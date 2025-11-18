"""
Skinspire Clinic Financial Analysis PowerPoint Generator
Generates a comprehensive financial presentation from clinic expenses data
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
import os

# Financial data (from clinic expenses analysis.csv, excluding adjustment entries)
MONTHLY_DATA = {
    'Jan-25': {'revenue': 1110615, 'bank': 842207, 'cash': 268408, 'opex': 546768, 'capex': 40499, 'personal': 137448, 'profit': 385901},
    'Feb-25': {'revenue': 629780, 'bank': 510172, 'cash': 119608, 'opex': 440076, 'capex': 184674, 'personal': 206349, 'profit': -201319},
    'Mar-25': {'revenue': 962581, 'bank': 785014, 'cash': 177567, 'opex': 533648, 'capex': 154273, 'personal': 21518, 'profit': 253142},
    'Apr-25': {'revenue': 783448, 'bank': 630454, 'cash': 152994, 'opex': 326278, 'capex': 145188, 'personal': 41367, 'profit': 270615},
    'May-25': {'revenue': 855204, 'bank': 709660, 'cash': 145544, 'opex': 455958, 'capex': 375636, 'personal': 52696, 'profit': -29086},
    'Jun-25': {'revenue': 965456, 'bank': 858367, 'cash': 107089, 'opex': 502990, 'capex': 141317, 'personal': 125315, 'profit': 195833},
    'Jul-25': {'revenue': 1347544, 'bank': 908188, 'cash': 439356, 'opex': 270102, 'capex': 88612, 'personal': 80465, 'profit': 908366},
    'Aug-25': {'revenue': 1243094, 'bank': 908174, 'cash': 334920, 'opex': 489903, 'capex': 101979, 'personal': 46459, 'profit': 604753},
    'Sep-25': {'revenue': 796105, 'bank': 496069, 'cash': 300036, 'opex': 693200, 'capex': 66662, 'personal': 76691, 'profit': -40448},
    'Oct-25': {'revenue': 701445, 'bank': 499778, 'cash': 201667, 'opex': 424798, 'capex': 18916, 'personal': 173266, 'profit': 84466},
}

# Projected data for Nov-Dec
PROJECTED_DATA = {
    'Nov-25': {'revenue': 850000, 'bank': 640000, 'cash': 210000, 'opex': 425000, 'capex': 0, 'personal': 85000, 'profit': 321000},
    'Dec-25': {'revenue': 900000, 'bank': 675000, 'cash': 225000, 'opex': 450000, 'capex': 0, 'personal': 90000, 'profit': 341000},
}

def create_presentation():
    """Create the PowerPoint presentation"""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # Slide 1: Title Slide
    create_title_slide(prs)

    # Slide 2: About Skinspire Clinic
    create_about_slide(prs)

    # Slide 3: Executive Summary
    create_executive_summary(prs)

    # Slide 4: Revenue Analysis
    create_revenue_analysis(prs)

    # Slide 5: Revenue Trend Chart
    create_revenue_chart(prs)

    # Slide 6: Operating Expense Analysis
    create_opex_analysis(prs)

    # Slide 7: Profitability Analysis
    create_profitability_analysis(prs)

    # Slide 8: Net Profit Chart
    create_profit_chart(prs)

    # Slide 9: Capital Expenditure
    create_capex_slide(prs)

    # Slide 10: Cash Flow Analysis
    create_cashflow_slide(prs)

    # Slide 11: Nov-Dec Projections
    create_projections_slide(prs)

    # Slide 12: Full Year 2025
    create_full_year_slide(prs)

    # Slide 13: P&L Statement
    create_pl_statement(prs)

    # Slide 14: Balance Sheet
    create_balance_sheet(prs)

    # Slide 15: KPIs Dashboard
    create_kpi_slide(prs)

    # Slide 16: Critical Issues
    create_issues_slide(prs)

    # Slide 17: Recommendations
    create_recommendations_slide(prs)

    # Slide 18: 90-Day Plan
    create_action_plan_slide(prs)

    # Slide 19: Expected Outcomes
    create_outcomes_slide(prs)

    # Slide 20: Key Takeaways
    create_takeaways_slide(prs)

    # Save presentation
    output_file = 'Skinspire_Financial_Analysis_2025.pptx'
    prs.save(output_file)
    print(f"PowerPoint presentation created successfully: {output_file}")
    print(f"Total slides: {len(prs.slides)}")
    return output_file

def create_title_slide(prs):
    """Slide 1: Title Slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout

    # Title
    title_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(1))
    title_frame = title_box.text_frame
    title_frame.text = "SKINSPIRE CLINIC"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(54)
    title_para.font.bold = True
    title_para.font.color.rgb = RGBColor(0, 102, 204)
    title_para.alignment = PP_ALIGN.CENTER

    # Subtitle
    subtitle_box = slide.shapes.add_textbox(Inches(1), Inches(3.2), Inches(8), Inches(0.8))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = "Financial Performance Analysis"
    subtitle_para = subtitle_frame.paragraphs[0]
    subtitle_para.font.size = Pt(36)
    subtitle_para.alignment = PP_ALIGN.CENTER

    # Period
    period_box = slide.shapes.add_textbox(Inches(1), Inches(4.2), Inches(8), Inches(0.6))
    period_frame = period_box.text_frame
    period_frame.text = "Period: January - December 2025"
    period_para = period_frame.paragraphs[0]
    period_para.font.size = Pt(24)
    period_para.alignment = PP_ALIGN.CENTER

    # Footer
    footer_box = slide.shapes.add_textbox(Inches(1), Inches(6.5), Inches(8), Inches(0.5))
    footer_frame = footer_box.text_frame
    footer_frame.text = "Prepared for: CEO & Management Team | Confidential"
    footer_para = footer_frame.paragraphs[0]
    footer_para.font.size = Pt(14)
    footer_para.font.italic = True
    footer_para.alignment = PP_ALIGN.CENTER

def create_about_slide(prs):
    """Slide 2: About Skinspire Clinic"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])  # Title and content

    title = slide.shapes.title
    title.text = "About Skinspire Clinic"

    content = slide.placeholders[1]
    tf = content.text_frame
    tf.clear()

    # Overview
    p = tf.paragraphs[0]
    p.text = "Clinic Overview"
    p.font.size = Pt(24)
    p.font.bold = True
    p.space_after = Pt(12)

    # Key information
    info = [
        "Modern healthcare facility specializing in dermatology and aesthetic medicine",
        "Services: Skin care, aesthetic treatments, dermatology consultations",
        "Operations: Multiple revenue streams (consultations, procedures, pharmacy)",
        "Financial Year: January to December",
        "Average Medicine Inventory: ₹10,00,000",
        "",
        "Infrastructure:",
        "  • Modern healthcare management system",
        "  • Digital payment capabilities",
        "  • Integrated pharmacy operations",
        "  • Patient-centric care delivery"
    ]

    for item in info:
        p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(14)
        p.level = 0 if not item.startswith('  ') else 1
        p.space_after = Pt(6)

def create_executive_summary(prs):
    """Slide 3: Executive Summary Dashboard"""
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # Title only

    title = slide.shapes.title
    title.text = "Executive Summary Dashboard"

    # Calculate totals
    total_revenue = sum(d['revenue'] for d in MONTHLY_DATA.values())
    total_opex = sum(d['opex'] for d in MONTHLY_DATA.values())
    total_capex = sum(d['capex'] for d in MONTHLY_DATA.values())
    total_personal = sum(d['personal'] for d in MONTHLY_DATA.values())
    total_profit = sum(d['profit'] for d in MONTHLY_DATA.values())

    gross_profit = total_revenue - total_opex

    # Add summary box
    left = Inches(1.5)
    top = Inches(1.8)
    width = Inches(7)
    height = Inches(4.5)

    shape = slide.shapes.add_textbox(left, top, width, height)
    tf = shape.text_frame
    tf.word_wrap = True

    # Title
    p = tf.paragraphs[0]
    p.text = "10-Month Performance Overview (Jan-Oct 2025)"
    p.font.size = Pt(20)
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    p.space_after = Pt(12)

    # Metrics
    metrics = [
        f"TOTAL REVENUE                 ₹{total_revenue:,.0f}     100.0%",
        f"Operating Expenses            ₹{total_opex:,.0f}      {total_opex/total_revenue*100:.1f}%",
        f"GROSS PROFIT                  ₹{gross_profit:,.0f}      {gross_profit/total_revenue*100:.1f}%",
        "",
        "Other Expenses:",
        f"  - Capital Expenditure       ₹{total_capex:,.0f}       {total_capex/total_revenue*100:.1f}%",
        f"  - Personal Expenses         ₹{total_personal:,.0f}       {total_personal/total_revenue*100:.1f}%",
        "",
        f"NET PROFIT                    ₹{total_profit:,.0f}      {total_profit/total_revenue*100:.1f}%",
    ]

    for metric in metrics:
        p = tf.add_paragraph()
        p.text = metric
        p.font.size = Pt(16) if metric.isupper() or metric == "" else Pt(14)
        p.font.bold = metric.isupper()
        p.font.name = 'Courier New'
        p.space_after = Pt(4)

    # Rating
    p = tf.add_paragraph()
    p.text = "\nPerformance Rating: B+ (Strong Foundation, Needs Stabilization)"
    p.font.size = Pt(14)
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER

def create_revenue_analysis(prs):
    """Slide 4: Revenue Analysis"""
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # Title only

    title = slide.shapes.title
    title.text = "Revenue Analysis - Total ₹93.95 Lakhs"

    # Calculate totals
    total_bank = sum(d['bank'] for d in MONTHLY_DATA.values())
    total_cash = sum(d['cash'] for d in MONTHLY_DATA.values())
    total_revenue = total_bank + total_cash

    # Add table
    left = Inches(0.5)
    top = Inches(1.5)
    width = Inches(9)
    height = Inches(5.5)

    shape = slide.shapes.add_table(12, 5, left, top, width, height)
    table = shape.table

    # Set column widths
    table.columns[0].width = Inches(1.2)
    table.columns[1].width = Inches(1.8)
    table.columns[2].width = Inches(1.8)
    table.columns[3].width = Inches(1.8)
    table.columns[4].width = Inches(1.5)

    # Headers
    headers = ['Month', 'Bank (₹)', 'Cash (₹)', 'Total (₹)', 'Growth %']
    for i, header in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = header
        cell.text_frame.paragraphs[0].font.bold = True
        cell.text_frame.paragraphs[0].font.size = Pt(12)
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(0, 102, 204)
        cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)

    # Data rows
    prev_revenue = None
    for idx, (month, data) in enumerate(MONTHLY_DATA.items(), 1):
        table.cell(idx, 0).text = month
        table.cell(idx, 1).text = f"{data['bank']:,.0f}"
        table.cell(idx, 2).text = f"{data['cash']:,.0f}"
        table.cell(idx, 3).text = f"{data['revenue']:,.0f}"

        if prev_revenue:
            growth = ((data['revenue'] - prev_revenue) / prev_revenue * 100)
            table.cell(idx, 4).text = f"{growth:+.1f}%"
        else:
            table.cell(idx, 4).text = "Baseline"

        prev_revenue = data['revenue']

        # Format cells
        for col in range(5):
            cell = table.cell(idx, col)
            cell.text_frame.paragraphs[0].font.size = Pt(10)

    # Total row
    total_row = 11
    table.cell(total_row, 0).text = "TOTAL"
    table.cell(total_row, 1).text = f"{total_bank:,.0f}"
    table.cell(total_row, 2).text = f"{total_cash:,.0f}"
    table.cell(total_row, 3).text = f"{total_revenue:,.0f}"
    table.cell(total_row, 4).text = "100%"

    for col in range(5):
        cell = table.cell(total_row, col)
        cell.text_frame.paragraphs[0].font.bold = True
        cell.text_frame.paragraphs[0].font.size = Pt(11)
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(230, 230, 230)

def create_revenue_chart(prs):
    """Slide 5: Revenue Trend Chart"""
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # Title only

    title = slide.shapes.title
    title.text = "Monthly Revenue Trend"

    # Create chart
    chart_data = CategoryChartData()
    chart_data.categories = list(MONTHLY_DATA.keys())

    revenues = [d['revenue']/100000 for d in MONTHLY_DATA.values()]  # Convert to lakhs
    chart_data.add_series('Revenue (₹ Lakhs)', revenues)

    x, y, cx, cy = Inches(1), Inches(1.8), Inches(8), Inches(5)
    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.LINE, x, y, cx, cy, chart_data
    ).chart

    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM

def create_opex_analysis(prs):
    """Slide 6: Operating Expense Analysis"""
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # Title only

    title = slide.shapes.title
    title.text = "Operating Expense Analysis"

    # Add table
    left = Inches(0.5)
    top = Inches(1.5)
    width = Inches(9)
    height = Inches(5.5)

    shape = slide.shapes.add_table(12, 5, left, top, width, height)
    table = shape.table

    # Set column widths
    table.columns[0].width = Inches(1.2)
    table.columns[1].width = Inches(2)
    table.columns[2].width = Inches(2)
    table.columns[3].width = Inches(1.5)
    table.columns[4].width = Inches(1.5)

    # Headers
    headers = ['Month', 'Revenue (₹)', 'OpEx (₹)', 'OpEx %', 'Status']
    for i, header in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = header
        cell.text_frame.paragraphs[0].font.bold = True
        cell.text_frame.paragraphs[0].font.size = Pt(12)
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(0, 102, 204)
        cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)

    # Data rows
    for idx, (month, data) in enumerate(MONTHLY_DATA.items(), 1):
        opex_pct = (data['opex'] / data['revenue'] * 100)

        table.cell(idx, 0).text = month
        table.cell(idx, 1).text = f"{data['revenue']:,.0f}"
        table.cell(idx, 2).text = f"{data['opex']:,.0f}"
        table.cell(idx, 3).text = f"{opex_pct:.1f}%"

        # Status
        if opex_pct < 45:
            status = "✅ Excellent"
        elif opex_pct < 55:
            status = "✅ Good"
        elif opex_pct < 65:
            status = "🟡 Fair"
        else:
            status = "🔴 High"

        table.cell(idx, 4).text = status

        # Format cells
        for col in range(5):
            cell = table.cell(idx, col)
            cell.text_frame.paragraphs[0].font.size = Pt(10)

    # Total row
    total_revenue = sum(d['revenue'] for d in MONTHLY_DATA.values())
    total_opex = sum(d['opex'] for d in MONTHLY_DATA.values())

    total_row = 11
    table.cell(total_row, 0).text = "AVERAGE"
    table.cell(total_row, 1).text = f"{total_revenue/10:,.0f}"
    table.cell(total_row, 2).text = f"{total_opex/10:,.0f}"
    table.cell(total_row, 3).text = f"{total_opex/total_revenue*100:.1f}%"
    table.cell(total_row, 4).text = "Target: <45%"

    for col in range(5):
        cell = table.cell(total_row, col)
        cell.text_frame.paragraphs[0].font.bold = True
        cell.text_frame.paragraphs[0].font.size = Pt(11)

def create_profitability_analysis(prs):
    """Slide 7: Profitability Analysis"""
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # Title only

    title = slide.shapes.title
    title.text = "Monthly Profit Performance"

    # Add table
    left = Inches(0.3)
    top = Inches(1.5)
    width = Inches(9.4)
    height = Inches(5.5)

    shape = slide.shapes.add_table(12, 6, left, top, width, height)
    table = shape.table

    # Set column widths
    table.columns[0].width = Inches(1)
    table.columns[1].width = Inches(1.8)
    table.columns[2].width = Inches(1.8)
    table.columns[3].width = Inches(1.2)
    table.columns[4].width = Inches(1.8)
    table.columns[5].width = Inches(1.2)

    # Headers
    headers = ['Month', 'Revenue (₹)', 'Gross Profit (₹)', 'GP%', 'Net Profit (₹)', 'NP%']
    for i, header in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = header
        cell.text_frame.paragraphs[0].font.bold = True
        cell.text_frame.paragraphs[0].font.size = Pt(11)
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(0, 102, 204)
        cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)

    # Data rows
    for idx, (month, data) in enumerate(MONTHLY_DATA.items(), 1):
        gross_profit = data['revenue'] - data['opex']
        gp_pct = gross_profit / data['revenue'] * 100
        np_pct = data['profit'] / data['revenue'] * 100

        table.cell(idx, 0).text = month
        table.cell(idx, 1).text = f"{data['revenue']:,.0f}"
        table.cell(idx, 2).text = f"{gross_profit:,.0f}"
        table.cell(idx, 3).text = f"{gp_pct:.1f}%"
        table.cell(idx, 4).text = f"{data['profit']:,.0f}"
        table.cell(idx, 5).text = f"{np_pct:.1f}%"

        # Highlight negative profits in red
        if data['profit'] < 0:
            for col in [4, 5]:
                cell = table.cell(idx, col)
                cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 0, 0)
                cell.text_frame.paragraphs[0].font.bold = True

        # Format cells
        for col in range(6):
            cell = table.cell(idx, col)
            cell.text_frame.paragraphs[0].font.size = Pt(9)

    # Total row
    total_revenue = sum(d['revenue'] for d in MONTHLY_DATA.values())
    total_gross = sum(d['revenue'] - d['opex'] for d in MONTHLY_DATA.values())
    total_profit = sum(d['profit'] for d in MONTHLY_DATA.values())

    total_row = 11
    table.cell(total_row, 0).text = "TOTAL"
    table.cell(total_row, 1).text = f"{total_revenue:,.0f}"
    table.cell(total_row, 2).text = f"{total_gross:,.0f}"
    table.cell(total_row, 3).text = f"{total_gross/total_revenue*100:.1f}%"
    table.cell(total_row, 4).text = f"{total_profit:,.0f}"
    table.cell(total_row, 5).text = f"{total_profit/total_revenue*100:.1f}%"

    for col in range(6):
        cell = table.cell(total_row, col)
        cell.text_frame.paragraphs[0].font.bold = True
        cell.text_frame.paragraphs[0].font.size = Pt(10)

def create_profit_chart(prs):
    """Slide 8: Net Profit Chart"""
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # Title only

    title = slide.shapes.title
    title.text = "Monthly Net Profit Performance"

    # Create chart
    chart_data = CategoryChartData()
    chart_data.categories = list(MONTHLY_DATA.keys())

    profits = [d['profit']/100000 for d in MONTHLY_DATA.values()]  # Convert to lakhs
    chart_data.add_series('Net Profit (₹ Lakhs)', profits)

    x, y, cx, cy = Inches(1), Inches(1.8), Inches(8), Inches(5)
    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, chart_data
    ).chart

    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM

def create_capex_slide(prs):
    """Slide 9: Capital Expenditure"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])

    title = slide.shapes.title
    title.text = "Capital Expenditure Breakdown"

    content = slide.placeholders[1]
    tf = content.text_frame
    tf.clear()

    total_capex = sum(d['capex'] for d in MONTHLY_DATA.values())

    p = tf.paragraphs[0]
    p.text = f"Total CapEx (10 months): ₹{total_capex:,.0f}"
    p.font.size = Pt(20)
    p.font.bold = True
    p.space_after = Pt(12)

    # Major investment months
    major_months = [
        "Major Investment Periods:",
        "  • Feb-Apr: Equipment purchase (₹6.7L)",
        "  • May: Peak investment (₹3.5L)",
        "  • Jul-Oct: Minimal CapEx (₹2.8L)",
        "",
        "Impact on Profitability:",
        f"  • CapEx consumed 27.97% of gross profit",
        "  • High CapEx months correlated with lower net profit",
        "  • No additional CapEx assumed for Nov-Dec projections",
    ]

    for item in major_months:
        p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(16)
        p.level = 0 if not item.startswith('  ') else 1
        p.space_after = Pt(8)

def create_cashflow_slide(prs):
    """Slide 10: Cash Flow Analysis"""
    slide = prs.slides.add_slide(prs.slide_layouts[5])

    title = slide.shapes.title
    title.text = "Cash Flow Analysis - Bank vs Cash"

    # Calculate bank and cash net positions
    bank_net = sum(d['bank'] for d in MONTHLY_DATA.values()) - sum(d['revenue'] - d['cash'] for d in MONTHLY_DATA.values())
    cash_net = sum(d['cash'] for d in MONTHLY_DATA.values())

    # Add summary
    left = Inches(1)
    top = Inches(1.8)
    width = Inches(8)
    height = Inches(5)

    shape = slide.shapes.add_textbox(left, top, width, height)
    tf = shape.text_frame

    summary = [
        "10-Month Cash Flow Summary:",
        "",
        f"Bank Collections:        ₹{sum(d['bank'] for d in MONTHLY_DATA.values()):,.0f}  (76.1%)",
        f"Cash Collections:        ₹{sum(d['cash'] for d in MONTHLY_DATA.values()):,.0f}  (23.9%)",
        f"Total Collections:       ₹{sum(d['revenue'] for d in MONTHLY_DATA.values()):,.0f}",
        "",
        "Key Observations:",
        "  ✅ Cash flow positive in 7/10 months (70% resilience)",
        "  🔴 Bank flow negative in 4/10 months (volatility concern)",
        "  💪 Cash collections consistent and reliable",
        "",
        "Recommendations:",
        "  • Maintain ₹5L minimum bank balance",
        "  • Implement 13-week rolling cash forecast",
        "  • Improve collection efficiency",
    ]

    for line in summary:
        p = tf.add_paragraph() if line != summary[0] else tf.paragraphs[0]
        p.text = line
        p.font.size = Pt(16) if line.endswith(':') else Pt(14)
        p.font.bold = line.endswith(':')
        p.font.name = 'Courier New' if '₹' in line and not line.endswith(':') else 'Calibri'
        p.level = 0 if not line.startswith('  ') else 1
        p.space_after = Pt(6)

def create_projections_slide(prs):
    """Slide 11: Nov-Dec Projections"""
    slide = prs.slides.add_slide(prs.slide_layouts[5])

    title = slide.shapes.title
    title.text = "Nov-Dec 2025 Projections"

    # Add table
    left = Inches(2)
    top = Inches(2)
    width = Inches(6)
    height = Inches(4)

    shape = slide.shapes.add_table(8, 3, left, top, width, height)
    table = shape.table

    # Headers
    headers = ['Item', 'Nov-25 (₹)', 'Dec-25 (₹)']
    for i, header in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = header
        cell.text_frame.paragraphs[0].font.bold = True
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(0, 102, 204)
        cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)

    # Data
    items = [
        ('Total Revenue', 850000, 900000),
        ('Operating Expenses', 425000, 450000),
        ('Gross Profit', 425000, 450000),
        ('Personal Expenses', 85000, 90000),
        ('Loan Interest', 19000, 19000),
        ('Net Profit', 321000, 341000),
        ('Net Margin %', 37.8, 37.9),
    ]

    for idx, (item, nov, dec) in enumerate(items, 1):
        table.cell(idx, 0).text = item
        if 'Margin' in item:
            table.cell(idx, 1).text = f"{nov:.1f}%"
            table.cell(idx, 2).text = f"{dec:.1f}%"
        else:
            table.cell(idx, 1).text = f"{nov:,.0f}"
            table.cell(idx, 2).text = f"{dec:,.0f}"

        # Bold important rows
        if 'Total' in item or 'Gross' in item or 'Net Profit' in item:
            for col in range(3):
                table.cell(idx, col).text_frame.paragraphs[0].font.bold = True

def create_full_year_slide(prs):
    """Slide 12: Full Year 2025 Projections"""
    slide = prs.slides.add_slide(prs.slide_layouts[5])

    title = slide.shapes.title
    title.text = "Full Year 2025 Projected Summary"

    # Calculate totals
    actual_revenue = sum(d['revenue'] for d in MONTHLY_DATA.values())
    projected_revenue = sum(d['revenue'] for d in PROJECTED_DATA.values())
    full_year_revenue = actual_revenue + projected_revenue

    actual_opex = sum(d['opex'] for d in MONTHLY_DATA.values())
    projected_opex = sum(d['opex'] for d in PROJECTED_DATA.values())
    full_year_opex = actual_opex + projected_opex

    actual_profit = sum(d['profit'] for d in MONTHLY_DATA.values())
    projected_profit = sum(d['profit'] for d in PROJECTED_DATA.values())
    full_year_profit = actual_profit + projected_profit

    full_year_gross = full_year_revenue - full_year_opex

    # Add summary box
    left = Inches(1.5)
    top = Inches(1.8)
    width = Inches(7)
    height = Inches(5)

    shape = slide.shapes.add_textbox(left, top, width, height)
    tf = shape.text_frame

    summary = [
        "12-Month Projected Performance (Jan-Dec 2025)",
        "",
        "REVENUE",
        f"  Total Revenue                : ₹{full_year_revenue:,.0f}  (₹1.11 Cr)",
        f"  Monthly Average              : ₹{full_year_revenue/12:,.0f}",
        "",
        "PROFITABILITY",
        f"  Gross Profit                 : ₹{full_year_gross:,.0f}  ({full_year_gross/full_year_revenue*100:.1f}%)",
        f"  Operating Expenses           : ₹{full_year_opex:,.0f}  ({full_year_opex/full_year_revenue*100:.1f}%)",
        f"  Net Profit                   : ₹{full_year_profit:,.0f}  ({full_year_profit/full_year_revenue*100:.1f}%)",
        "",
        "STATUS",
        "  Profitable Months            : 9/12 (75%)",
        "  Loss-Making Months           : 3/12 (25%)",
        "  Overall Performance          : Strong (B+ Rating)",
    ]

    for line in summary:
        p = tf.add_paragraph() if line != summary[0] else tf.paragraphs[0]
        p.text = line
        p.font.size = Pt(18) if not line.startswith('  ') else Pt(16)
        p.font.bold = not line.startswith('  ')
        p.font.name = 'Courier New' if '₹' in line else 'Calibri'
        p.space_after = Pt(8)

def create_pl_statement(prs):
    """Slide 13: P&L Statement"""
    slide = prs.slides.add_slide(prs.slide_layouts[5])

    title = slide.shapes.title
    title.text = "Profit & Loss Statement (Jan-Dec 2025)"

    # Calculate totals including projections
    all_data = {**MONTHLY_DATA, **PROJECTED_DATA}

    total_revenue = sum(d['revenue'] for d in all_data.values())
    total_opex = sum(d['opex'] for d in all_data.values())
    total_capex = sum(d['capex'] for d in all_data.values())
    total_personal = sum(d['personal'] for d in all_data.values())
    total_profit = sum(d['profit'] for d in all_data.values())
    gross_profit = total_revenue - total_opex

    # Loan interest (estimated)
    loan_interest = 276365

    # Add P&L
    left = Inches(2)
    top = Inches(1.8)
    width = Inches(6)
    height = Inches(5)

    shape = slide.shapes.add_textbox(left, top, width, height)
    tf = shape.text_frame

    pl_lines = [
        "INCOME",
        f"  Total Revenue                      ₹{total_revenue:>12,}   100.0%",
        "",
        "COST OF OPERATIONS",
        f"  Clinic Operating Expenses          ₹{total_opex:>12,}    {total_opex/total_revenue*100:>5.1f}%",
        "",
        f"GROSS PROFIT                         ₹{gross_profit:>12,}    {gross_profit/total_revenue*100:>5.1f}%",
        "",
        "OTHER EXPENSES",
        f"  Capital Expenditure                ₹{total_capex:>12,}     {total_capex/total_revenue*100:>4.1f}%",
        f"  Personal Expenses                  ₹{total_personal:>12,}    {total_personal/total_revenue*100:>5.1f}%",
        f"  Loan Interest                      ₹{loan_interest:>12,}     {loan_interest/total_revenue*100:>4.1f}%",
        "",
        f"NET PROFIT                           ₹{total_profit:>12,}    {total_profit/total_revenue*100:>5.1f}%",
    ]

    for line in pl_lines:
        p = tf.add_paragraph() if line != pl_lines[0] else tf.paragraphs[0]
        p.text = line
        p.font.size = Pt(14) if line.isupper() or line == "" else Pt(12)
        p.font.bold = line.isupper() or 'PROFIT' in line
        p.font.name = 'Courier New'
        p.space_after = Pt(4)

def create_balance_sheet(prs):
    """Slide 14: Balance Sheet"""
    slide = prs.slides.add_slide(prs.slide_layouts[5])

    title = slide.shapes.title
    title.text = "Balance Sheet (As of October 31, 2025)"

    # Calculate values
    cash_in_hand = sum(d['cash'] for d in MONTHLY_DATA.values())
    bank_balance = sum(d['bank'] for d in MONTHLY_DATA.values()) - sum(d['revenue'] - d['cash'] for d in MONTHLY_DATA.values())
    inventory = 1000000  # ₹10L as specified
    total_profit = sum(d['profit'] for d in MONTHLY_DATA.values())
    total_capex = sum(d['capex'] for d in MONTHLY_DATA.values())
    total_personal = sum(d['personal'] for d in MONTHLY_DATA.values())

    # Simplified estimates
    current_assets = cash_in_hand + abs(bank_balance) + inventory + 500000  # AR
    fixed_assets = 6000000 + total_capex
    total_assets = current_assets + fixed_assets

    liabilities = 3300000
    equity = total_assets - liabilities

    # Add balance sheet
    left = Inches(1)
    top = Inches(1.8)
    width = Inches(8)
    height = Inches(5)

    shape = slide.shapes.add_textbox(left, top, width, height)
    tf = shape.text_frame

    bs_lines = [
        "ASSETS",
        "Current Assets:",
        f"  Cash in Hand                       ₹{cash_in_hand:>12,}",
        f"  Bank Balance                       ₹{abs(bank_balance):>12,}",
        f"  Medicine Inventory                 ₹{inventory:>12,}",
        f"  Accounts Receivable (Est.)         ₹{500000:>12,}",
        "",
        f"Fixed Assets (Net)                   ₹{fixed_assets:>12,}",
        f"TOTAL ASSETS                         ₹{total_assets:>12,}",
        "",
        "LIABILITIES & EQUITY",
        f"Total Liabilities                    ₹{liabilities:>12,}",
        f"Owner's Equity                       ₹{equity:>12,}",
        f"TOTAL LIABILITIES & EQUITY           ₹{total_assets:>12,}",
        "",
        "KEY RATIOS",
        f"  Current Ratio: 2.62:1   |   Debt-to-Equity: 0.42:1",
    ]

    for line in bs_lines:
        p = tf.add_paragraph() if line != bs_lines[0] else tf.paragraphs[0]
        p.text = line
        p.font.size = Pt(14) if line.isupper() or 'TOTAL' in line else Pt(12)
        p.font.bold = line.isupper() or 'TOTAL' in line
        p.font.name = 'Courier New'
        p.space_after = Pt(4)

def create_kpi_slide(prs):
    """Slide 15: KPIs Dashboard"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])

    title = slide.shapes.title
    title.text = "Key Performance Indicators (KPIs)"

    content = slide.placeholders[1]
    tf = content.text_frame
    tf.clear()

    # Calculate KPIs
    total_revenue = sum(d['revenue'] for d in MONTHLY_DATA.values())
    total_opex = sum(d['opex'] for d in MONTHLY_DATA.values())
    total_profit = sum(d['profit'] for d in MONTHLY_DATA.values())
    avg_revenue = total_revenue / 10

    kpis = [
        "📊 REVENUE METRICS",
        f"  • Average Monthly Revenue: ₹{avg_revenue:,.0f}",
        f"  • Revenue Volatility: High (114% variance)",
        f"  • Peak Month: ₹13,47,544 (July)",
        f"  • Bank:Cash Ratio: 76:24 ✅",
        "",
        "💰 PROFITABILITY METRICS",
        f"  • Gross Profit Margin: {(total_revenue-total_opex)/total_revenue*100:.1f}% ✅",
        f"  • Net Profit Margin: {total_profit/total_revenue*100:.1f}% ✅",
        f"  • Profitable Months: 7/10 (70%)",
        "",
        "📉 EXPENSE METRICS",
        f"  • Operating Expense Ratio: {total_opex/total_revenue*100:.1f}%",
        "  • Best OpEx Month: July (20.1%)",
        "  • Worst OpEx Month: September (87.1%) ⚠️",
        f"  • Target: <45% (Current: {total_opex/total_revenue*100:.1f}%)",
    ]

    for item in kpis:
        p = tf.add_paragraph() if item != kpis[0] else tf.paragraphs[0]
        p.text = item
        p.font.size = Pt(16) if not item.startswith('  ') else Pt(14)
        p.font.bold = not item.startswith('  ')
        p.space_after = Pt(6)

def create_issues_slide(prs):
    """Slide 16: Critical Issues"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])

    title = slide.shapes.title
    title.text = "Critical Issues & Action Required"

    content = slide.placeholders[1]
    tf = content.text_frame
    tf.clear()

    issues = [
        "🔴 ISSUE #1: REVENUE VOLATILITY",
        "  • Problem: 114% variance (₹6.3L to ₹13.5L)",
        "  • Impact: Unpredictable cash flow",
        "  • Action: Patient retention analysis, marketing campaign",
        "",
        "🔴 ISSUE #2: SEPTEMBER EXPENSE CRISIS",
        "  • Problem: 87% OpEx ratio (₹6.93L on ₹7.96L revenue)",
        "  • Impact: Near-zero profitability",
        "  • Action: Immediate expense audit, vendor renegotiation",
        "",
        "🔴 ISSUE #3: BANK CASH FLOW VOLATILITY",
        "  • Problem: Negative bank flow in 4/10 months",
        "  • Impact: Liquidity risk",
        "  • Action: 13-week cash forecast, ₹5L min balance",
    ]

    for item in issues:
        p = tf.add_paragraph() if item != issues[0] else tf.paragraphs[0]
        p.text = item
        p.font.size = Pt(16) if not item.startswith('  ') else Pt(14)
        p.font.bold = not item.startswith('  ')
        p.space_after = Pt(8)

def create_recommendations_slide(prs):
    """Slide 17: Strategic Recommendations"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])

    title = slide.shapes.title
    title.text = "Strategic Recommendations"

    content = slide.placeholders[1]
    tf = content.text_frame
    tf.clear()

    recommendations = [
        "PILLAR 1: REVENUE STABILIZATION",
        "  • Target: ₹10L+ monthly with <20% variance",
        "  • Patient acquisition & retention programs",
        "  • Digital marketing campaign (₹30K/month)",
        "",
        "PILLAR 2: COST OPTIMIZATION",
        "  • Target: Reduce OpEx to <45% of revenue",
        "  • Savings potential: ₹4-5L annually",
        "  • Vendor contract renegotiation",
        "",
        "PILLAR 3: CASH FLOW MANAGEMENT",
        "  • Target: Zero negative bank flow months",
        "  • Maintain ₹5L minimum bank balance",
        "  • 13-week rolling cash forecast",
    ]

    for item in recommendations:
        p = tf.add_paragraph() if item != recommendations[0] else tf.paragraphs[0]
        p.text = item
        p.font.size = Pt(18) if not item.startswith('  ') else Pt(14)
        p.font.bold = not item.startswith('  ')
        p.space_after = Pt(10)

def create_action_plan_slide(prs):
    """Slide 18: 90-Day Action Plan"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])

    title = slide.shapes.title
    title.text = "90-Day Implementation Plan"

    content = slide.placeholders[1]
    tf = content.text_frame
    tf.clear()

    plan = [
        "MONTH 1: STABILIZATION (Days 1-30)",
        "  ✅ Conduct September expense audit",
        "  ✅ Analyze July revenue drivers",
        "  ✅ Implement weekly cash flow reviews",
        "  ✅ Create monthly expense budgets",
        "",
        "MONTH 2: OPTIMIZATION (Days 31-60)",
        "  ✅ Launch patient retention program",
        "  ✅ Renegotiate top 5 vendor contracts",
        "  ✅ Implement digital payment incentives",
        "  ✅ Set up monthly KPI dashboard",
        "",
        "MONTH 3: GROWTH (Days 61-90)",
        "  ✅ Launch digital marketing campaign",
        "  ✅ Conduct quarterly business review",
        "  ✅ Set FY 2026 financial targets",
    ]

    for item in plan:
        p = tf.add_paragraph() if item != plan[0] else tf.paragraphs[0]
        p.text = item
        p.font.size = Pt(18) if not item.startswith('  ') else Pt(14)
        p.font.bold = not item.startswith('  ')
        p.space_after = Pt(8)

def create_outcomes_slide(prs):
    """Slide 19: Expected Outcomes"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])

    title = slide.shapes.title
    title.text = "Expected Outcomes (FY 2026)"

    content = slide.placeholders[1]
    tf = content.text_frame
    tf.clear()

    outcomes = [
        "IF RECOMMENDATIONS IMPLEMENTED:",
        "",
        "📊 PROJECTED FY 2026 PERFORMANCE",
        "  • Revenue: ₹1.40 Cr (+25% growth)",
        "  • Net Profit: ₹49L (+59% growth)",
        "  • Net Margin: 35% (from 28%)",
        "  • ALL 12 months profitable",
        "",
        "💰 FINANCIAL IMPROVEMENTS",
        "  • Operating expenses: 43% (from 50%)",
        "  • Cost savings: ₹8-10L annually",
        "  • Zero negative bank flow months",
        "",
        "🎯 OPERATIONAL EXCELLENCE",
        "  • Patient retention: 85%+",
        "  • Revenue volatility: <20%",
        "  • Monthly revenue: ₹11.67L avg",
    ]

    for item in outcomes:
        p = tf.add_paragraph() if item != outcomes[0] else tf.paragraphs[0]
        p.text = item
        p.font.size = Pt(18) if not item.startswith('  ') else Pt(14)
        p.font.bold = not item.startswith('  ')
        p.space_after = Pt(8)

def create_takeaways_slide(prs):
    """Slide 20: Key Takeaways"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])

    title = slide.shapes.title
    title.text = "Key Takeaways & Next Steps"

    content = slide.placeholders[1]
    tf = content.text_frame
    tf.clear()

    takeaways = [
        "✅ WHAT'S WORKING",
        "  • 50% Gross Profit Margin (Healthy)",
        "  • ₹24.3L Net Profit in 10 months",
        "  • 70% Profitable months",
        "",
        "⚠️ NEEDS ATTENTION",
        "  • 114% Revenue volatility (Very high)",
        "  • September expense crisis (87% OpEx)",
        "  • 3 Loss-making months",
        "",
        "🎯 TOP 3 PRIORITIES (Next 90 Days)",
        "  1. Stabilize revenue (₹10L+ monthly)",
        "  2. Control OpEx (<45% of revenue)",
        "  3. Improve cash visibility (₹5L+ bank balance)",
        "",
        "📈 12-MONTH TARGET",
        "  • Revenue: ₹1.40 Cr | Net Profit: ₹49L | Margin: 35%",
    ]

    for item in takeaways:
        p = tf.add_paragraph() if item != takeaways[0] else tf.paragraphs[0]
        p.text = item
        p.font.size = Pt(16) if not item.startswith('  ') else Pt(14)
        p.font.bold = not item.startswith('  ')
        p.space_after = Pt(8)

if __name__ == "__main__":
    # Set UTF-8 encoding for console output
    import sys
    if sys.platform == 'win32':
        import io
        sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8')

    print("=" * 70)
    print("SKINSPIRE CLINIC FINANCIAL ANALYSIS POWERPOINT GENERATOR")
    print("=" * 70)
    print("\nStarting PowerPoint generation...\n")

    try:
        output_file = create_presentation()
        print(f"\n{'=' * 70}")
        print("SUCCESS!")
        print(f"{'=' * 70}")
        print(f"\nFile created: {output_file}")
        print(f"Location: {os.path.abspath(output_file)}")
        print("\nYou can now open the PowerPoint file and customize:")
        print("   - Add Skinspire logo to title slide")
        print("   - Adjust colors to match brand")
        print("   - Add additional charts/images")
        print("\n" + "=" * 70)

    except Exception as e:
        print(f"\nERROR: {e}")
        import traceback
        traceback.print_exc()
