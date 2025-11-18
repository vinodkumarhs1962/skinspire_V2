"""
Enhanced update script that also updates chart data
Preserves all formatting, logos, and custom changes
"""

from pptx import Presentation
from pptx.chart.data import CategoryChartData
import os

# Corrected financial data
MONTHLY_DATA = {
    'Jan': {'revenue': 1110615, 'opex': 546768, 'capex': 40499, 'personal': 137448, 'profit': 385901},
    'Feb': {'revenue': 629780, 'opex': 440076, 'capex': 354674, 'personal': 206349, 'profit': -371319},
    'Mar': {'revenue': 962581, 'opex': 533648, 'capex': 404273, 'personal': 21518, 'profit': 3142},
    'Apr': {'revenue': 783448, 'opex': 326278, 'capex': 395188, 'personal': 41367, 'profit': 20615},
    'May': {'revenue': 855204, 'opex': 455958, 'capex': 405636, 'personal': 52696, 'profit': -59086},
    'Jun': {'revenue': 965456, 'opex': 502990, 'capex': 141317, 'personal': 125315, 'profit': 195833},
    'Jul': {'revenue': 1347544, 'opex': 270102, 'capex': 157612, 'personal': 80465, 'profit': 839366},
    'Aug': {'revenue': 1243094, 'opex': 489903, 'capex': 101979, 'personal': 46459, 'profit': 604753},
    'Sep': {'revenue': 796105, 'opex': 693200, 'capex': 66662, 'personal': 76691, 'profit': -40448},
    'Oct': {'revenue': 701445, 'opex': 424798, 'capex': 18916, 'personal': 173266, 'profit': 84466},
}

def update_text_in_shape(shape, replacements):
    """Update text in a shape"""
    if not shape.has_text_frame:
        return 0

    updates = 0
    text_frame = shape.text_frame
    for paragraph in text_frame.paragraphs:
        for run in paragraph.runs:
            for old_text, new_text in replacements.items():
                if old_text in run.text:
                    run.text = run.text.replace(old_text, new_text)
                    updates += 1
    return updates

def update_text_in_table(table, replacements):
    """Update text in table cells"""
    updates = 0
    for row in table.rows:
        for cell in row.cells:
            for old_text, new_text in replacements.items():
                if old_text in cell.text:
                    cell.text = cell.text.replace(old_text, new_text)
                    updates += 1
    return updates

def try_update_chart(shape, slide_num):
    """Attempt to update chart data"""
    if not shape.has_chart:
        return False

    try:
        chart = shape.chart
        chart_data = chart.chart_data

        # Try to identify chart type by categories
        if hasattr(chart_data, 'categories'):
            categories = [str(cat) for cat in chart_data.categories]

            # Monthly revenue/profit charts
            if 'Jan' in categories and 'Oct' in categories:
                print(f"    Found monthly chart with {len(categories)} categories")

                # Update chart data
                new_chart_data = CategoryChartData()
                new_chart_data.categories = categories

                # Determine what data to use based on number of series
                num_series = len(chart_data.series)

                if num_series == 1:
                    series_name = chart_data.series[0].name

                    if 'Revenue' in str(series_name):
                        values = [MONTHLY_DATA[m]['revenue']/100000 for m in categories]
                        new_chart_data.add_series(series_name, values)
                        print(f"      Updated Revenue series")
                    elif 'Profit' in str(series_name):
                        values = [MONTHLY_DATA[m]['profit']/100000 for m in categories]
                        new_chart_data.add_series(series_name, values)
                        print(f"      Updated Profit series")
                    elif 'OpEx' in str(series_name) or 'Expense' in str(series_name):
                        values = [(MONTHLY_DATA[m]['opex']/MONTHLY_DATA[m]['revenue']*100) for m in categories]
                        new_chart_data.add_series(series_name, values)
                        print(f"      Updated OpEx % series")

                elif num_series == 2:
                    # Likely Gross Profit vs Net Profit
                    for series in chart_data.series:
                        if 'Gross' in str(series.name):
                            gross_values = [(MONTHLY_DATA[m]['revenue'] - MONTHLY_DATA[m]['opex'])/100000 for m in categories]
                            new_chart_data.add_series(series.name, gross_values)
                            print(f"      Updated Gross Profit series")
                        elif 'Net' in str(series.name):
                            net_values = [MONTHLY_DATA[m]['profit']/100000 for m in categories]
                            new_chart_data.add_series(series.name, net_values)
                            print(f"      Updated Net Profit series")

                # Replace chart data
                chart.replace_data(new_chart_data)
                return True

    except Exception as e:
        print(f"    Could not auto-update chart: {e}")
        return False

    return False

def update_presentation(input_file, output_file):
    """Update data in existing presentation"""
    print(f"\nOpening: {input_file}")
    prs = Presentation(input_file)

    total_updates = 0
    charts_updated = 0

    # Define all replacements
    replacements = {
        # Net profit (10-month)
        '24,32,223': '16,63,223',
        '24.32L': '16.63L',
        '₹24.32L': '₹16.63L',

        # Net margin (10-month)
        '25.89%': '17.70%',
        '25.9%': '17.7%',

        # Total CapEx
        '13,17,756': '20,86,756',
        '13.18L': '20.87L',
        '₹13.18L': '₹20.87L',

        # February
        '-2,01,319': '-3,71,319',
        '-201319': '-371319',
        '-2.01L': '-3.71L',
        '₹-2.01L': '₹-3.71L',

        # March
        '2,53,142': '3,142',
        '253142': '3142',
        '2.53L': '0.03L',
        '₹2.53L': '₹0.03L',

        # April
        '2,70,615': '20,615',
        '270615': '20615',
        '2.71L': '0.21L',
        '₹2.71L': '₹0.21L',

        # May
        '-29,086': '-59,086',
        '-0.29L': '-0.59L',
        '₹-0.29L': '₹-0.59L',

        # July
        '9,08,366': '8,39,366',
        '908366': '839366',
        '9.08L': '8.39L',
        '₹9.08L': '₹8.39L',

        # CapEx months
        '1,84,674': '3,54,674',
        '184674': '354674',
        '1,54,273': '4,04,273',
        '154273': '404273',
        '1,45,188': '3,95,188',
        '145188': '395188',
        '3,75,636': '4,05,636',
        '375636': '405636',
        '88,612': '1,57,612',

        # Profitable months
        '7/10': '5/10',

        # CapEx as % of GP
        '27.97%': '44.29%',

        # Full year
        '30,94,223': '23,25,223',
        '30.94L': '23.25L',
        '₹30.94L': '₹23.25L',
        '27.8%': '20.9%',
    }

    # Process each slide
    for slide_num, slide in enumerate(prs.slides, 1):
        print(f"\nSlide {slide_num}:")
        slide_updates = 0

        for shape in slide.shapes:
            # Update text
            updates = update_text_in_shape(shape, replacements)
            if updates > 0:
                print(f"  Text updates: {updates}")
                slide_updates += updates

            # Update tables
            if shape.has_table:
                updates = update_text_in_table(shape.table, replacements)
                if updates > 0:
                    print(f"  Table updates: {updates}")
                    slide_updates += updates

            # Update charts
            if shape.has_chart:
                if try_update_chart(shape, slide_num):
                    charts_updated += 1
                    slide_updates += 1

        total_updates += slide_updates
        if slide_updates > 0:
            print(f"  Total slide updates: {slide_updates}")

    # Save
    print(f"\nSaving: {output_file}")
    prs.save(output_file)

    print("\n" + "=" * 70)
    print("✅ UPDATE COMPLETE!")
    print("=" * 70)
    print(f"Text/Table updates: {total_updates - charts_updated}")
    print(f"Charts updated: {charts_updated}")
    print(f"Total slides: {len(prs.slides)}")
    print(f"\nOutput file: {os.path.abspath(output_file)}")

    if charts_updated < 4:
        print("\n⚠️  NOTE: Some charts may need manual verification")
        print("   Please review all charts to ensure data is correct")

if __name__ == "__main__":
    import sys

    # Set UTF-8 encoding
    if sys.platform == 'win32':
        import io
        sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8')

    print("=" * 70)
    print("SKINSPIRE CLINIC - ENHANCED DATA UPDATE")
    print("Updates text, tables, AND charts")
    print("=" * 70)

    input_file = 'Skinspire_Financial_Analysis_2025_Formatted.pptx'
    output_file = 'Skinspire_Financial_Analysis_2025_Final.pptx'

    if not os.path.exists(input_file):
        print(f"\n❌ ERROR: File not found: {input_file}")
        print("\nPlease ensure your edited PowerPoint file is named:")
        print(f"  {input_file}")
        sys.exit(1)

    try:
        update_presentation(input_file, output_file)
        print("\n✅ Your formatting, logos, and styles are preserved!")
        print("=" * 70)
    except Exception as e:
        print(f"\n❌ ERROR: {e}")
        import traceback
        traceback.print_exc()
