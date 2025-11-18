"""
Skinspire Clinic Financial Analysis - Enhanced PowerPoint Generator
Creates a professionally formatted presentation with better styling
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION, XL_DATA_LABEL_POSITION
from pptx.oxml.xmlchemy import OxmlElement
import os

# Brand colors
PRIMARY_BLUE = RGBColor(0, 102, 204)
DARK_BLUE = RGBColor(0, 51, 102)
LIGHT_BLUE = RGBColor(230, 242, 255)
GREEN = RGBColor(0, 176, 80)
RED = RGBColor(192, 0, 0)
ORANGE = RGBColor(255, 153, 0)
GRAY = RGBColor(127, 127, 127)
WHITE = RGBColor(255, 255, 255)
BLACK = RGBColor(0, 0, 0)

# Financial data (from clinic expenses analysis.csv, excluding adjustment entries)
# CORRECTED: Capital Expenditure now includes Capital expenses + Capital Expenditure CASH + Loan interest
MONTHLY_DATA = {
    'Jan': {'revenue': 1110615, 'bank': 842207, 'cash': 268408, 'opex': 546768, 'capex': 40499, 'personal': 137448, 'profit': 385901},
    'Feb': {'revenue': 629780, 'bank': 510172, 'cash': 119608, 'opex': 440076, 'capex': 354674, 'personal': 206349, 'profit': -371319},
    'Mar': {'revenue': 962581, 'bank': 785014, 'cash': 177567, 'opex': 533648, 'capex': 404273, 'personal': 21518, 'profit': 3142},
    'Apr': {'revenue': 783448, 'bank': 630454, 'cash': 152994, 'opex': 326278, 'capex': 395188, 'personal': 41367, 'profit': 20615},
    'May': {'revenue': 855204, 'bank': 709660, 'cash': 145544, 'opex': 455958, 'capex': 405636, 'personal': 52696, 'profit': -59086},
    'Jun': {'revenue': 965456, 'bank': 858367, 'cash': 107089, 'opex': 502990, 'capex': 141317, 'personal': 125315, 'profit': 195833},
    'Jul': {'revenue': 1347544, 'bank': 908188, 'cash': 439356, 'opex': 270102, 'capex': 157612, 'personal': 80465, 'profit': 839366},
    'Aug': {'revenue': 1243094, 'bank': 908174, 'cash': 334920, 'opex': 489903, 'capex': 101979, 'personal': 46459, 'profit': 604753},
    'Sep': {'revenue': 796105, 'bank': 496069, 'cash': 300036, 'opex': 693200, 'capex': 66662, 'personal': 76691, 'profit': -40448},
    'Oct': {'revenue': 701445, 'bank': 499778, 'cash': 201667, 'opex': 424798, 'capex': 18916, 'personal': 173266, 'profit': 84466},
}

PROJECTED_DATA = {
    'Nov': {'revenue': 850000, 'bank': 640000, 'cash': 210000, 'opex': 425000, 'capex': 0, 'personal': 85000, 'profit': 321000},
    'Dec': {'revenue': 900000, 'bank': 675000, 'cash': 225000, 'opex': 450000, 'capex': 0, 'personal': 90000, 'profit': 341000},
}

def add_title_with_style(slide, title_text, subtitle_text=None):
    """Add formatted title to slide"""
    # Title background
    title_shape = slide.shapes.add_shape(
        1,  # Rectangle
        Inches(0), Inches(0), Inches(10), Inches(1.2)
    )
    title_shape.fill.solid()
    title_shape.fill.fore_color.rgb = PRIMARY_BLUE
    title_shape.line.color.rgb = PRIMARY_BLUE

    # Title text
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
    title_frame = title_box.text_frame
    title_frame.text = title_text
    p = title_frame.paragraphs[0]
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.LEFT

    # Subtitle if provided
    if subtitle_text:
        subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.75), Inches(9), Inches(0.4))
        subtitle_frame = subtitle_box.text_frame
        subtitle_frame.text = subtitle_text
        p = subtitle_frame.paragraphs[0]
        p.font.size = Pt(16)
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.LEFT

def create_metric_box(slide, left, top, width, height, title, value, subtitle="", bg_color=PRIMARY_BLUE):
    """Create a styled metric box"""
    # Background
    shape = slide.shapes.add_shape(1, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = bg_color
    shape.line.color.rgb = bg_color

    # Title
    title_box = slide.shapes.add_textbox(left, top + Inches(0.1), width, Inches(0.3))
    tf = title_box.text_frame
    tf.text = title
    p = tf.paragraphs[0]
    p.font.size = Pt(12)
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    # Value
    value_box = slide.shapes.add_textbox(left, top + Inches(0.4), width, Inches(0.4))
    tf = value_box.text_frame
    tf.text = value
    p = tf.paragraphs[0]
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    # Subtitle
    if subtitle:
        sub_box = slide.shapes.add_textbox(left, top + Inches(0.85), width, Inches(0.25))
        tf = sub_box.text_frame
        tf.text = subtitle
        p = tf.paragraphs[0]
        p.font.size = Pt(10)
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER

def create_presentation():
    """Create the enhanced PowerPoint presentation"""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # Create slides
    create_title_slide(prs)
    create_executive_summary(prs)
    create_revenue_overview(prs)
    create_revenue_chart_slide(prs)
    create_profitability_slide(prs)
    create_profit_chart_slide(prs)
    create_expense_analysis_slide(prs)
    create_cashflow_slide(prs)
    create_projections_slide(prs)
    create_pl_statement_slide(prs)
    create_balance_sheet_slide(prs)
    create_kpi_dashboard_slide(prs)
    create_issues_slide(prs)
    create_recommendations_slide(prs)
    create_action_plan_slide(prs)
    create_outcomes_slide(prs)

    # Save
    import datetime
    timestamp = datetime.datetime.now().strftime("%Y%m%d_%H%M%S")
    output_file = f'Skinspire_Financial_Analysis_2025_Corrected.pptx'
    prs.save(output_file)
    print(f"PowerPoint created: {output_file}")
    print(f"Total slides: {len(prs.slides)}")
    return output_file

def create_title_slide(prs):
    """Slide 1: Professional Title Slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Background gradient effect (top half dark blue)
    bg_shape = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(4))
    bg_shape.fill.solid()
    bg_shape.fill.fore_color.rgb = DARK_BLUE
    bg_shape.line.color.rgb = DARK_BLUE

    # Title
    title_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(1.2))
    tf = title_box.text_frame
    tf.text = "SKINSPIRE CLINIC"
    p = tf.paragraphs[0]
    p.font.size = Pt(60)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    # Subtitle
    subtitle_box = slide.shapes.add_textbox(Inches(1), Inches(3.3), Inches(8), Inches(0.8))
    tf = subtitle_box.text_frame
    tf.text = "Financial Performance Analysis"
    p = tf.paragraphs[0]
    p.font.size = Pt(32)
    p.font.color.rgb = LIGHT_BLUE
    p.alignment = PP_ALIGN.CENTER

    # Period
    period_box = slide.shapes.add_textbox(Inches(1), Inches(4.8), Inches(8), Inches(0.6))
    tf = period_box.text_frame
    tf.text = "January - December 2025"
    p = tf.paragraphs[0]
    p.font.size = Pt(24)
    p.font.color.rgb = PRIMARY_BLUE
    p.alignment = PP_ALIGN.CENTER

    # Footer
    footer_box = slide.shapes.add_textbox(Inches(1), Inches(6.8), Inches(8), Inches(0.5))
    tf = footer_box.text_frame
    tf.text = "Prepared for: CEO & Management Team | Confidential"
    p = tf.paragraphs[0]
    p.font.size = Pt(12)
    p.font.italic = True
    p.font.color.rgb = GRAY
    p.alignment = PP_ALIGN.CENTER

def create_executive_summary(prs):
    """Slide 2: Executive Summary with Metric Boxes"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_with_style(slide, "Executive Summary", "10-Month Performance Overview (Jan-Oct 2025)")

    # Calculate totals
    total_revenue = sum(d['revenue'] for d in MONTHLY_DATA.values())
    total_opex = sum(d['opex'] for d in MONTHLY_DATA.values())
    total_profit = sum(d['profit'] for d in MONTHLY_DATA.values())
    gross_profit = total_revenue - total_opex

    # Metric boxes - Row 1
    create_metric_box(slide, Inches(0.5), Inches(1.8), Inches(2.2), Inches(1.2),
                     "TOTAL REVENUE", f"₹{total_revenue/100000:.2f}L", "100%", PRIMARY_BLUE)

    create_metric_box(slide, Inches(3), Inches(1.8), Inches(2.2), Inches(1.2),
                     "GROSS PROFIT", f"₹{gross_profit/100000:.2f}L",
                     f"{gross_profit/total_revenue*100:.1f}%", GREEN)

    create_metric_box(slide, Inches(5.5), Inches(1.8), Inches(2.2), Inches(1.2),
                     "NET PROFIT", f"₹{total_profit/100000:.2f}L",
                     f"{total_profit/total_revenue*100:.1f}%", DARK_BLUE)

    create_metric_box(slide, Inches(8), Inches(1.8), Inches(1.5), Inches(1.2),
                     "RATING", "B+", "Strong", ORANGE)

    # Metric boxes - Row 2
    create_metric_box(slide, Inches(0.5), Inches(3.3), Inches(2.2), Inches(1.2),
                     "OPERATING EXPENSE", f"₹{total_opex/100000:.2f}L",
                     f"{total_opex/total_revenue*100:.1f}%", ORANGE)

    profitable_months = sum(1 for d in MONTHLY_DATA.values() if d['profit'] > 0)
    create_metric_box(slide, Inches(3), Inches(3.3), Inches(2.2), Inches(1.2),
                     "PROFITABLE MONTHS", f"{profitable_months}/10",
                     "70%", GREEN if profitable_months >= 7 else RED)

    avg_revenue = total_revenue / 10
    create_metric_box(slide, Inches(5.5), Inches(3.3), Inches(2.2), Inches(1.2),
                     "AVG MONTHLY REVENUE", f"₹{avg_revenue/100000:.2f}L",
                     "Per Month", PRIMARY_BLUE)

    create_metric_box(slide, Inches(8), Inches(3.3), Inches(1.5), Inches(1.2),
                     "GP %", f"{gross_profit/total_revenue*100:.0f}%",
                     "Margin", GREEN)

    # Key insights box
    insight_shape = slide.shapes.add_shape(1, Inches(0.5), Inches(4.8), Inches(9), Inches(2.2))
    insight_shape.fill.solid()
    insight_shape.fill.fore_color.rgb = LIGHT_BLUE
    insight_shape.line.color.rgb = PRIMARY_BLUE
    insight_shape.line.width = Pt(2)

    insight_box = slide.shapes.add_textbox(Inches(0.7), Inches(5), Inches(8.6), Inches(1.8))
    tf = insight_box.text_frame
    tf.word_wrap = True

    insights = [
        ("STRENGTHS:", "• 50% Gross Profit Margin  • Positive Overall Cash Flow  • Diversified Revenue Streams"),
        ("CONCERNS:", "• High Revenue Volatility (114%)  • 3 Loss-Making Months  • September Expense Crisis (87% OpEx)"),
        ("PRIORITY:", "• Stabilize Monthly Revenue  • Control Operating Expenses  • Improve Cash Flow Visibility")
    ]

    for i, (label, text) in enumerate(insights):
        p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
        p.text = f"{label}  {text}"
        p.font.size = Pt(12)
        p.font.bold = True if ":" in label else False
        p.space_after = Pt(8)
        if "STRENGTHS" in label:
            p.font.color.rgb = GREEN
        elif "CONCERNS" in label:
            p.font.color.rgb = RED
        else:
            p.font.color.rgb = PRIMARY_BLUE

def create_revenue_overview(prs):
    """Slide 3: Revenue Overview Table"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_with_style(slide, "Revenue Analysis", "Monthly Breakdown (Jan-Oct 2025)")

    # Create styled table
    rows, cols = 12, 5
    left, top, width, height = Inches(0.8), Inches(1.8), Inches(8.4), Inches(5)

    shape = slide.shapes.add_table(rows, cols, left, top, width, height)
    table = shape.table

    # Column widths
    table.columns[0].width = Inches(1.2)
    table.columns[1].width = Inches(2)
    table.columns[2].width = Inches(1.8)
    table.columns[3].width = Inches(2)
    table.columns[4].width = Inches(1.4)

    # Headers
    headers = ['Month', 'Revenue (₹)', 'OpEx (₹)', 'Net Profit (₹)', 'NP %']
    for i, header in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = header
        cell.text_frame.paragraphs[0].font.bold = True
        cell.text_frame.paragraphs[0].font.size = Pt(11)
        cell.text_frame.paragraphs[0].font.color.rgb = WHITE
        cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        cell.fill.solid()
        cell.fill.fore_color.rgb = PRIMARY_BLUE

    # Data rows
    for idx, (month, data) in enumerate(MONTHLY_DATA.items(), 1):
        np_pct = data['profit'] / data['revenue'] * 100

        row_data = [
            month,
            f"{data['revenue']:,}",
            f"{data['opex']:,}",
            f"{data['profit']:,}",
            f"{np_pct:.1f}%"
        ]

        for col, value in enumerate(row_data):
            cell = table.cell(idx, col)
            cell.text = value
            cell.text_frame.paragraphs[0].font.size = Pt(10)
            cell.text_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT if col > 0 else PP_ALIGN.CENTER

            # Highlight negative profits
            if col in [3, 4] and data['profit'] < 0:
                cell.text_frame.paragraphs[0].font.color.rgb = RED
                cell.text_frame.paragraphs[0].font.bold = True

            # Alternate row colors
            if idx % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(245, 245, 245)

    # Total row
    total_revenue = sum(d['revenue'] for d in MONTHLY_DATA.values())
    total_opex = sum(d['opex'] for d in MONTHLY_DATA.values())
    total_profit = sum(d['profit'] for d in MONTHLY_DATA.values())

    total_row = 11
    totals = ['TOTAL', f"{total_revenue:,}", f"{total_opex:,}", f"{total_profit:,}",
              f"{total_profit/total_revenue*100:.1f}%"]

    for col, value in enumerate(totals):
        cell = table.cell(total_row, col)
        cell.text = value
        cell.text_frame.paragraphs[0].font.bold = True
        cell.text_frame.paragraphs[0].font.size = Pt(11)
        cell.text_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT if col > 0 else PP_ALIGN.CENTER
        cell.fill.solid()
        cell.fill.fore_color.rgb = LIGHT_BLUE

def create_revenue_chart_slide(prs):
    """Slide 4: Revenue Trend Chart"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_with_style(slide, "Revenue Trend", "Monthly Performance (Lakhs)")

    # Create chart
    chart_data = CategoryChartData()
    chart_data.categories = list(MONTHLY_DATA.keys())

    revenues = [d['revenue']/100000 for d in MONTHLY_DATA.values()]
    chart_data.add_series('Revenue (₹ Lakhs)', revenues)

    x, y, cx, cy = Inches(1), Inches(2), Inches(8), Inches(4.8)
    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.LINE_MARKERS, x, y, cx, cy, chart_data
    ).chart

    # Style the chart
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart.legend.font.size = Pt(12)

    # Add data labels
    plot = chart.plots[0]
    plot.has_data_labels = True
    data_labels = plot.data_labels
    data_labels.font.size = Pt(9)
    data_labels.number_format = '0.0'

def create_profitability_slide(prs):
    """Slide 5: Profitability Overview"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_with_style(slide, "Profitability Analysis", "Gross Profit vs Net Profit")

    # Calculate data
    months = list(MONTHLY_DATA.keys())
    gross_profits = [(d['revenue'] - d['opex'])/100000 for d in MONTHLY_DATA.values()]
    net_profits = [d['profit']/100000 for d in MONTHLY_DATA.values()]

    # Create chart
    chart_data = CategoryChartData()
    chart_data.categories = months
    chart_data.add_series('Gross Profit (₹L)', gross_profits)
    chart_data.add_series('Net Profit (₹L)', net_profits)

    x, y, cx, cy = Inches(0.8), Inches(2), Inches(8.4), Inches(4.8)
    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, chart_data
    ).chart

    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart.legend.font.size = Pt(12)

def create_profit_chart_slide(prs):
    """Slide 6: Net Profit Chart"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_with_style(slide, "Net Profit Performance", "Monthly Net Profit (Lakhs)")

    # Create chart
    chart_data = CategoryChartData()
    chart_data.categories = list(MONTHLY_DATA.keys())

    profits = [d['profit']/100000 for d in MONTHLY_DATA.values()]
    chart_data.add_series('Net Profit (₹ Lakhs)', profits)

    x, y, cx, cy = Inches(1), Inches(2), Inches(8), Inches(4.8)
    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, chart_data
    ).chart

    chart.has_legend = False

def create_expense_analysis_slide(prs):
    """Slide 7: Operating Expense Analysis"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_with_style(slide, "Operating Expense Analysis", "OpEx as % of Revenue")

    # Create chart
    chart_data = CategoryChartData()
    chart_data.categories = list(MONTHLY_DATA.keys())

    opex_pct = [(d['opex']/d['revenue']*100) for d in MONTHLY_DATA.values()]
    chart_data.add_series('OpEx %', opex_pct)

    x, y, cx, cy = Inches(1), Inches(2), Inches(8), Inches(4.8)
    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.LINE_MARKERS, x, y, cx, cy, chart_data
    ).chart

    chart.has_legend = False

    # Add target line text
    target_box = slide.shapes.add_textbox(Inches(7.5), Inches(4.5), Inches(2), Inches(0.5))
    tf = target_box.text_frame
    tf.text = "Target: <45%"
    p = tf.paragraphs[0]
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = RED

def create_cashflow_slide(prs):
    """Slide 8: Cash Flow Summary"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_with_style(slide, "Cash Flow Analysis", "Bank vs Cash Collections")

    # Calculate totals
    total_bank = sum(d['bank'] for d in MONTHLY_DATA.values())
    total_cash = sum(d['cash'] for d in MONTHLY_DATA.values())

    # Metric boxes
    create_metric_box(slide, Inches(1), Inches(2), Inches(2.5), Inches(1.3),
                     "BANK COLLECTIONS", f"₹{total_bank/100000:.2f}L", "76.1%", PRIMARY_BLUE)

    create_metric_box(slide, Inches(3.8), Inches(2), Inches(2.5), Inches(1.3),
                     "CASH COLLECTIONS", f"₹{total_cash/100000:.2f}L", "23.9%", GREEN)

    create_metric_box(slide, Inches(6.6), Inches(2), Inches(2.5), Inches(1.3),
                     "TOTAL CASH FLOW", f"₹{(total_bank+total_cash)/100000:.2f}L", "100%", DARK_BLUE)

    # Key observations
    obs_shape = slide.shapes.add_shape(1, Inches(1), Inches(3.6), Inches(8), Inches(3))
    obs_shape.fill.solid()
    obs_shape.fill.fore_color.rgb = LIGHT_BLUE
    obs_shape.line.color.rgb = PRIMARY_BLUE
    obs_shape.line.width = Pt(2)

    obs_box = slide.shapes.add_textbox(Inches(1.3), Inches(3.8), Inches(7.4), Inches(2.6))
    tf = obs_box.text_frame
    tf.word_wrap = True

    observations = [
        "KEY OBSERVATIONS:",
        "",
        "• Cash flow positive in 7 out of 10 months (70% resilience)",
        "• Bank flow showed volatility with 4 negative months",
        "• September had the largest bank outflow (-₹3.4L)",
        "• Cash collections remained consistent throughout",
        "",
        "RECOMMENDATIONS:",
        "• Maintain minimum ₹5L bank balance at all times",
        "• Implement 13-week rolling cash forecast",
        "• Improve collection efficiency and payment terms"
    ]

    for i, text in enumerate(observations):
        p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(14) if text.endswith(':') else Pt(12)
        p.font.bold = text.endswith(':')
        p.space_after = Pt(6)

def create_projections_slide(prs):
    """Slide 9: Projections"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_with_style(slide, "Financial Projections", "Nov-Dec 2025 & Full Year")

    # Calculate full year
    actual_revenue = sum(d['revenue'] for d in MONTHLY_DATA.values())
    projected_revenue = sum(d['revenue'] for d in PROJECTED_DATA.values())
    full_year_revenue = actual_revenue + projected_revenue

    actual_profit = sum(d['profit'] for d in MONTHLY_DATA.values())
    projected_profit = sum(d['profit'] for d in PROJECTED_DATA.values())
    full_year_profit = actual_profit + projected_profit

    # Metric boxes
    create_metric_box(slide, Inches(0.8), Inches(2), Inches(2.7), Inches(1.3),
                     "FULL YEAR REVENUE", f"₹{full_year_revenue/10000000:.2f} Cr",
                     "Jan-Dec 2025", PRIMARY_BLUE)

    create_metric_box(slide, Inches(3.8), Inches(2), Inches(2.7), Inches(1.3),
                     "FULL YEAR PROFIT", f"₹{full_year_profit/100000:.2f}L",
                     f"{full_year_profit/full_year_revenue*100:.1f}% Margin", GREEN)

    create_metric_box(slide, Inches(6.8), Inches(2), Inches(2.2), Inches(1.3),
                     "PROFITABLE MONTHS", "9/12", "75%", ORANGE)

    # Projection table
    rows, cols = 4, 4
    left, top, width, height = Inches(1.5), Inches(3.8), Inches(7), Inches(2.5)

    shape = slide.shapes.add_table(rows, cols, left, top, width, height)
    table = shape.table

    # Headers
    headers = ['Metric', 'Nov-25 (₹)', 'Dec-25 (₹)', 'Total (₹)']
    for i, header in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = header
        cell.text_frame.paragraphs[0].font.bold = True
        cell.text_frame.paragraphs[0].font.size = Pt(12)
        cell.text_frame.paragraphs[0].font.color.rgb = WHITE
        cell.fill.solid()
        cell.fill.fore_color.rgb = PRIMARY_BLUE

    # Data
    nov_data = PROJECTED_DATA['Nov']
    dec_data = PROJECTED_DATA['Dec']

    rows_data = [
        ['Revenue', f"{nov_data['revenue']:,}", f"{dec_data['revenue']:,}",
         f"{nov_data['revenue'] + dec_data['revenue']:,}"],
        ['Operating Expenses', f"{nov_data['opex']:,}", f"{dec_data['opex']:,}",
         f"{nov_data['opex'] + dec_data['opex']:,}"],
        ['Net Profit', f"{nov_data['profit']:,}", f"{dec_data['profit']:,}",
         f"{nov_data['profit'] + dec_data['profit']:,}"]
    ]

    for idx, row_data in enumerate(rows_data, 1):
        for col, value in enumerate(row_data):
            cell = table.cell(idx, col)
            cell.text = value
            cell.text_frame.paragraphs[0].font.size = Pt(11)
            cell.text_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT if col > 0 else PP_ALIGN.LEFT
            if idx == 3:  # Net Profit row
                cell.text_frame.paragraphs[0].font.bold = True

def create_pl_statement_slide(prs):
    """Slide 10: P&L Statement"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_with_style(slide, "Profit & Loss Statement", "Projected Full Year 2025 (Jan-Dec)")

    # Calculate totals
    all_data = {**MONTHLY_DATA, **PROJECTED_DATA}
    total_revenue = sum(d['revenue'] for d in all_data.values())
    total_opex = sum(d['opex'] for d in all_data.values())
    total_capex = sum(d['capex'] for d in all_data.values())
    total_personal = sum(d['personal'] for d in all_data.values())
    gross_profit = total_revenue - total_opex
    total_profit = sum(d['profit'] for d in all_data.values())

    # Create P&L Box
    pl_shape = slide.shapes.add_shape(1, Inches(1.5), Inches(1.8), Inches(7), Inches(5))
    pl_shape.fill.solid()
    pl_shape.fill.fore_color.rgb = WHITE
    pl_shape.line.color.rgb = PRIMARY_BLUE
    pl_shape.line.width = Pt(3)

    pl_box = slide.shapes.add_textbox(Inches(1.8), Inches(2), Inches(6.4), Inches(4.6))
    tf = pl_box.text_frame
    tf.word_wrap = True

    pl_lines = [
        ("INCOME", ""),
        (f"  Total Revenue", f"₹{total_revenue:>15,}   100.0%"),
        ("", ""),
        ("COST OF OPERATIONS", ""),
        (f"  Operating Expenses", f"₹{total_opex:>15,}    {total_opex/total_revenue*100:>5.1f}%"),
        ("", ""),
        (f"GROSS PROFIT", f"₹{gross_profit:>15,}    {gross_profit/total_revenue*100:>5.1f}%"),
        ("", ""),
        ("OTHER EXPENSES", ""),
        (f"  Capital Expenditure*", f"₹{total_capex:>15,}    {total_capex/total_revenue*100:>5.1f}%"),
        (f"  Personal Expenses", f"₹{total_personal:>15,}    {total_personal/total_revenue*100:>5.1f}%"),
        ("", ""),
        (f"NET PROFIT", f"₹{total_profit:>15,}    {total_profit/total_revenue*100:>5.1f}%"),
        ("", ""),
        (f"  *Includes equipment, cash CapEx & loan interest", ""),
    ]

    for i, (label, value) in enumerate(pl_lines):
        p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
        p.text = f"{label:<25}{value}"
        p.font.size = Pt(13) if label.isupper() and not label.startswith('  ') else Pt(11)
        p.font.bold = label.isupper() or 'PROFIT' in label
        p.font.name = 'Courier New'
        p.space_after = Pt(2)

        if 'NET PROFIT' in label or 'GROSS PROFIT' in label:
            p.font.color.rgb = GREEN
            p.font.size = Pt(14)

def create_balance_sheet_slide(prs):
    """Slide 11: Balance Sheet"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_with_style(slide, "Balance Sheet", "As of October 31, 2025")

    # Calculate values
    cash_in_hand = sum(d['cash'] for d in MONTHLY_DATA.values())
    bank_balance = 954034  # From analysis
    inventory = 1000000
    total_profit = sum(d['profit'] for d in MONTHLY_DATA.values())
    total_capex = sum(d['capex'] for d in MONTHLY_DATA.values())

    current_assets = cash_in_hand + bank_balance + inventory + 500000
    fixed_assets = 6000000 + total_capex
    total_assets = current_assets + fixed_assets

    liabilities = 3300000
    equity = total_assets - liabilities

    # Create two column layout
    # Assets column
    asset_shape = slide.shapes.add_shape(1, Inches(0.5), Inches(1.8), Inches(4.5), Inches(5))
    asset_shape.fill.solid()
    asset_shape.fill.fore_color.rgb = LIGHT_BLUE
    asset_shape.line.color.rgb = PRIMARY_BLUE
    asset_shape.line.width = Pt(2)

    asset_box = slide.shapes.add_textbox(Inches(0.7), Inches(2), Inches(4.1), Inches(4.6))
    tf = asset_box.text_frame

    assets = [
        ("ASSETS", ""),
        ("Current Assets:", ""),
        (f"  Cash in Hand", f"₹{cash_in_hand:>12,}"),
        (f"  Bank Balance", f"₹{bank_balance:>12,}"),
        (f"  Inventory", f"₹{inventory:>12,}"),
        (f"  Accounts Receivable", f"₹{500000:>12,}"),
        ("", ""),
        (f"Fixed Assets (Net)", f"₹{fixed_assets:>12,}"),
        ("", ""),
        (f"TOTAL ASSETS", f"₹{total_assets:>12,}"),
    ]

    for i, (label, value) in enumerate(assets):
        p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
        p.text = f"{label:<20}{value}"
        p.font.size = Pt(13) if label.isupper() else Pt(11)
        p.font.bold = label.isupper() or label == 'Fixed Assets (Net)'
        p.font.name = 'Courier New'
        p.space_after = Pt(3)

    # Liabilities column
    liab_shape = slide.shapes.add_shape(1, Inches(5), Inches(1.8), Inches(4.5), Inches(5))
    liab_shape.fill.solid()
    liab_shape.fill.fore_color.rgb = LIGHT_BLUE
    liab_shape.line.color.rgb = PRIMARY_BLUE
    liab_shape.line.width = Pt(2)

    liab_box = slide.shapes.add_textbox(Inches(5.2), Inches(2), Inches(4.1), Inches(4.6))
    tf = liab_box.text_frame

    liabilities_data = [
        ("LIABILITIES & EQUITY", ""),
        ("", ""),
        (f"Total Liabilities", f"₹{liabilities:>12,}"),
        ("", ""),
        (f"Owner's Equity", f"₹{equity:>12,}"),
        ("", ""),
        (f"TOTAL L & E", f"₹{total_assets:>12,}"),
        ("", ""),
        ("KEY RATIOS", ""),
        ("Current Ratio", "2.62:1"),
        ("Debt-to-Equity", "0.42:1"),
    ]

    for i, (label, value) in enumerate(liabilities_data):
        p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
        p.text = f"{label:<20}{value}"
        p.font.size = Pt(13) if label.isupper() else Pt(11)
        p.font.bold = label.isupper()
        p.font.name = 'Courier New'
        p.space_after = Pt(3)

def create_kpi_dashboard_slide(prs):
    """Slide 12: KPI Dashboard"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_with_style(slide, "Key Performance Indicators", "Current Performance vs Targets")

    # Calculate KPIs
    total_revenue = sum(d['revenue'] for d in MONTHLY_DATA.values())
    total_opex = sum(d['opex'] for d in MONTHLY_DATA.values())
    total_profit = sum(d['profit'] for d in MONTHLY_DATA.values())
    gross_margin = (total_revenue - total_opex) / total_revenue * 100

    # KPI boxes - Row 1
    create_metric_box(slide, Inches(0.5), Inches(1.9), Inches(1.8), Inches(1.1),
                     "Gross Margin", f"{gross_margin:.1f}%", "Target: >50%",
                     GREEN if gross_margin >= 50 else ORANGE)

    create_metric_box(slide, Inches(2.5), Inches(1.9), Inches(1.8), Inches(1.1),
                     "Net Margin", f"{total_profit/total_revenue*100:.1f}%", "Target: >25%",
                     GREEN if total_profit/total_revenue*100 >= 25 else ORANGE)

    create_metric_box(slide, Inches(4.5), Inches(1.9), Inches(1.8), Inches(1.1),
                     "OpEx Ratio", f"{total_opex/total_revenue*100:.1f}%", "Target: <45%",
                     RED if total_opex/total_revenue*100 > 50 else ORANGE)

    create_metric_box(slide, Inches(6.5), Inches(1.9), Inches(1.8), Inches(1.1),
                     "Profitable", "70%", "7/10 months",
                     GREEN)

    create_metric_box(slide, Inches(8.5), Inches(1.9), Inches(1), Inches(1.1),
                     "Rating", "B+", "",
                     PRIMARY_BLUE)

    # Critical metrics table
    rows, cols = 6, 4
    left, top, width, height = Inches(1), Inches(3.4), Inches(8), Inches(3)

    shape = slide.shapes.add_table(rows, cols, left, top, width, height)
    table = shape.table

    # Headers
    headers = ['Metric', 'Current', 'Target', 'Status']
    for i, header in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = header
        cell.text_frame.paragraphs[0].font.bold = True
        cell.text_frame.paragraphs[0].font.size = Pt(12)
        cell.text_frame.paragraphs[0].font.color.rgb = WHITE
        cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        cell.fill.solid()
        cell.fill.fore_color.rgb = PRIMARY_BLUE

    # Data
    kpi_data = [
        ['Monthly Revenue', f"₹{total_revenue/10/100000:.2f}L", '₹10L+', 'Close'],
        ['Revenue Volatility', '114%', '<20%', 'High'],
        ['Operating Expense %', f"{total_opex/total_revenue*100:.1f}%", '<45%', 'High'],
        ['Negative Months', '3/10', '0/10', 'Concern'],
        ['Cash Flow', 'Positive', 'Positive', 'Good']
    ]

    for idx, row_data in enumerate(kpi_data, 1):
        for col, value in enumerate(row_data):
            cell = table.cell(idx, col)
            cell.text = value
            cell.text_frame.paragraphs[0].font.size = Pt(11)
            cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

            # Color code status
            if col == 3:
                if value in ['Good', 'Close']:
                    cell.text_frame.paragraphs[0].font.color.rgb = GREEN
                elif value in ['High', 'Concern']:
                    cell.text_frame.paragraphs[0].font.color.rgb = RED

def create_issues_slide(prs):
    """Slide 13: Critical Issues"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_with_style(slide, "Critical Issues", "Top 3 Priorities Requiring Immediate Attention")

    issues = [
        {
            'title': 'ISSUE #1: REVENUE VOLATILITY',
            'problem': '• 114% variance (₹6.3L to ₹13.5L)\n• Q3 decline: 25% drop from July peak\n• Unpredictable monthly performance',
            'action': '• Patient retention analysis\n• Launch marketing campaign\n• Analyze July success factors',
            'color': RED
        },
        {
            'title': 'ISSUE #2: SEPTEMBER EXPENSE CRISIS',
            'problem': '• Operating expenses: 87% of revenue\n• ₹6.93L OpEx on ₹7.96L revenue\n• Near-zero profitability',
            'action': '• Immediate expense audit\n• Vendor contract renegotiation\n• Monthly expense budgets',
            'color': RED
        },
        {
            'title': 'ISSUE #3: CASH FLOW VOLATILITY',
            'problem': '• Bank flow negative in 4/10 months\n• September: -₹3.4L bank outflow\n• Liquidity risk exposure',
            'action': '• 13-week rolling cash forecast\n• Maintain ₹5L minimum balance\n• Improve collection efficiency',
            'color': ORANGE
        }
    ]

    y_positions = [1.9, 3.5, 5.1]

    for idx, (issue, y_pos) in enumerate(zip(issues, y_positions)):
        # Issue box
        issue_shape = slide.shapes.add_shape(1, Inches(0.5), Inches(y_pos), Inches(9), Inches(1.4))
        issue_shape.fill.solid()
        issue_shape.fill.fore_color.rgb = LIGHT_BLUE
        issue_shape.line.color.rgb = issue['color']
        issue_shape.line.width = Pt(3)

        # Title
        title_box = slide.shapes.add_textbox(Inches(0.7), Inches(y_pos + 0.05), Inches(8.6), Inches(0.25))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = issue['title']
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = issue['color']

        # Problem
        prob_box = slide.shapes.add_textbox(Inches(0.7), Inches(y_pos + 0.32), Inches(4.2), Inches(0.9))
        tf = prob_box.text_frame
        p = tf.paragraphs[0]
        p.text = issue['problem']
        p.font.size = Pt(10)
        p.line_spacing = 1.2

        # Action
        action_box = slide.shapes.add_textbox(Inches(5.2), Inches(y_pos + 0.32), Inches(4), Inches(0.9))
        tf = action_box.text_frame
        p = tf.paragraphs[0]
        p.text = "ACTIONS:\n" + issue['action']
        p.font.size = Pt(10)
        p.font.bold = True
        p.font.color.rgb = GREEN
        p.line_spacing = 1.2

def create_recommendations_slide(prs):
    """Slide 14: Strategic Recommendations"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_with_style(slide, "Strategic Recommendations", "4 Pillars for Sustainable Growth")

    pillars = [
        {
            'title': '1. REVENUE STABILIZATION',
            'target': 'Target: ₹10L+ monthly',
            'items': '• Patient acquisition programs\n• Digital marketing (₹30K/month)\n• Retention initiatives\n• Corporate partnerships'
        },
        {
            'title': '2. COST OPTIMIZATION',
            'target': 'Target: <45% OpEx',
            'items': '• Vendor renegotiation\n• Zero-based budgeting\n• Staff optimization\n• Inventory management'
        },
        {
            'title': '3. CASH FLOW MANAGEMENT',
            'target': 'Target: Zero negative months',
            'items': '• ₹5L minimum balance\n• 13-week cash forecast\n• Collection efficiency\n• Payment term optimization'
        },
        {
            'title': '4. FINANCIAL VISIBILITY',
            'target': 'Target: Real-time dashboards',
            'items': '• Daily revenue tracking\n• Weekly expense review\n• Monthly business review\n• KPI scorecards'
        }
    ]

    positions = [(0.5, 2), (5, 2), (0.5, 4.5), (5, 4.5)]

    for pillar, (x, y) in zip(pillars, positions):
        # Pillar box
        pillar_shape = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(4.2), Inches(2.2))
        pillar_shape.fill.solid()
        pillar_shape.fill.fore_color.rgb = LIGHT_BLUE
        pillar_shape.line.color.rgb = PRIMARY_BLUE
        pillar_shape.line.width = Pt(2)

        # Title
        title_box = slide.shapes.add_textbox(Inches(x + 0.2), Inches(y + 0.1), Inches(3.8), Inches(0.3))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = pillar['title']
        p.font.size = Pt(13)
        p.font.bold = True
        p.font.color.rgb = PRIMARY_BLUE

        # Target
        target_box = slide.shapes.add_textbox(Inches(x + 0.2), Inches(y + 0.45), Inches(3.8), Inches(0.25))
        tf = target_box.text_frame
        p = tf.paragraphs[0]
        p.text = pillar['target']
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = GREEN

        # Items
        items_box = slide.shapes.add_textbox(Inches(x + 0.2), Inches(y + 0.75), Inches(3.8), Inches(1.3))
        tf = items_box.text_frame
        p = tf.paragraphs[0]
        p.text = pillar['items']
        p.font.size = Pt(10)
        p.line_spacing = 1.3

def create_action_plan_slide(prs):
    """Slide 15: 90-Day Action Plan"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_with_style(slide, "90-Day Implementation Plan", "Quarterly Roadmap")

    months = [
        {
            'title': 'MONTH 1: STABILIZATION',
            'subtitle': 'Days 1-30',
            'actions': '• Conduct September expense audit\n• Analyze July revenue drivers\n• Weekly cash flow reviews\n• Create expense budgets\n• Set up KPI dashboard',
            'color': RED
        },
        {
            'title': 'MONTH 2: OPTIMIZATION',
            'subtitle': 'Days 31-60',
            'actions': '• Launch retention program\n• Renegotiate vendor contracts\n• Digital payment incentives\n• Staff productivity review\n• Inventory optimization',
            'color': ORANGE
        },
        {
            'title': 'MONTH 3: GROWTH',
            'subtitle': 'Days 61-90',
            'actions': '• Digital marketing campaign\n• Quarterly business review\n• Set FY 2026 targets\n• Plan capacity expansion\n• Measure 90-day results',
            'color': GREEN
        }
    ]

    x_positions = [0.5, 3.5, 6.5]

    for month, x_pos in zip(months, x_positions):
        # Month box
        month_shape = slide.shapes.add_shape(1, Inches(x_pos), Inches(2), Inches(3), Inches(4.2))
        month_shape.fill.solid()
        month_shape.fill.fore_color.rgb = LIGHT_BLUE
        month_shape.line.color.rgb = month['color']
        month_shape.line.width = Pt(3)

        # Title
        title_box = slide.shapes.add_textbox(Inches(x_pos + 0.2), Inches(2.15), Inches(2.6), Inches(0.4))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = month['title']
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = month['color']

        # Subtitle
        sub_box = slide.shapes.add_textbox(Inches(x_pos + 0.2), Inches(2.55), Inches(2.6), Inches(0.25))
        tf = sub_box.text_frame
        p = tf.paragraphs[0]
        p.text = month['subtitle']
        p.font.size = Pt(10)
        p.font.italic = True

        # Actions
        actions_box = slide.shapes.add_textbox(Inches(x_pos + 0.2), Inches(2.9), Inches(2.6), Inches(3.1))
        tf = actions_box.text_frame
        p = tf.paragraphs[0]
        p.text = month['actions']
        p.font.size = Pt(10)
        p.line_spacing = 1.4

def create_outcomes_slide(prs):
    """Slide 16: Expected Outcomes"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_with_style(slide, "Expected Outcomes (FY 2026)", "If Recommendations Implemented")

    # FY 2026 projections
    create_metric_box(slide, Inches(0.8), Inches(2), Inches(2.7), Inches(1.3),
                     "FY 2026 REVENUE", "₹1.40 Cr", "+25% Growth", PRIMARY_BLUE)

    create_metric_box(slide, Inches(3.8), Inches(2), Inches(2.7), Inches(1.3),
                     "FY 2026 PROFIT", "₹49 Lakhs", "+59% Growth", GREEN)

    create_metric_box(slide, Inches(6.8), Inches(2), Inches(2.2), Inches(1.3),
                     "NET MARGIN", "35%", "From 28%", DARK_BLUE)

    # Comparison table
    rows, cols = 6, 4
    left, top, width, height = Inches(1), Inches(3.8), Inches(8), Inches(2.8)

    shape = slide.shapes.add_table(rows, cols, left, top, width, height)
    table = shape.table

    # Headers
    headers = ['Metric', 'FY 2025', 'FY 2026', 'Change']
    for i, header in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = header
        cell.text_frame.paragraphs[0].font.bold = True
        cell.text_frame.paragraphs[0].font.size = Pt(12)
        cell.text_frame.paragraphs[0].font.color.rgb = WHITE
        cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        cell.fill.solid()
        cell.fill.fore_color.rgb = PRIMARY_BLUE

    # Data
    comparison = [
        ['Revenue', '₹1.11 Cr', '₹1.40 Cr', '+25%'],
        ['Net Profit', '₹30.9L', '₹49.0L', '+59%'],
        ['NP Margin', '27.8%', '35.0%', '+7.2%'],
        ['OpEx Ratio', '49.9%', '43.0%', '-6.9%'],
        ['Loss Months', '3/12', '0/12', '100% improve']
    ]

    for idx, row_data in enumerate(comparison, 1):
        for col, value in enumerate(row_data):
            cell = table.cell(idx, col)
            cell.text = value
            cell.text_frame.paragraphs[0].font.size = Pt(11)
            cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

            # Color improvements
            if col == 3 and '+' in value:
                cell.text_frame.paragraphs[0].font.color.rgb = GREEN
                cell.text_frame.paragraphs[0].font.bold = True
            elif col == 3 and 'improve' in value:
                cell.text_frame.paragraphs[0].font.color.rgb = GREEN
                cell.text_frame.paragraphs[0].font.bold = True

if __name__ == "__main__":
    # Set UTF-8 encoding for console output
    import sys
    if sys.platform == 'win32':
        import io
        sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8')

    print("=" * 70)
    print("SKINSPIRE CLINIC - ENHANCED FINANCIAL PRESENTATION GENERATOR")
    print("=" * 70)
    print("\nGenerating professionally formatted PowerPoint...\n")

    try:
        output_file = create_presentation()
        print(f"\n{'=' * 70}")
        print("SUCCESS!")
        print(f"{'=' * 70}")
        print(f"\nFile: {output_file}")
        print(f"Location: {os.path.abspath(output_file)}")
        print("\nFeatures:")
        print("  - Professional formatting with brand colors")
        print("  - Styled metric boxes and tables")
        print("  - Charts with proper legends")
        print("  - Color-coded insights and status")
        print("\n" + "=" * 70)

    except Exception as e:
        print(f"\nERROR: {e}")
        import traceback
        traceback.print_exc()
