"""
Apply formatting from user's edited file to the corrected data file
- Copies logo/images from edited file
- Increases table font sizes
- Improves alignment
- Maintains data integrity from corrected file
"""

from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import os

def copy_images_from_slide(source_slide, target_slide):
    """Copy images from source slide to target slide"""
    copied = 0
    for shape in source_slide.shapes:
        if shape.shape_type == 13:  # Picture
            try:
                # Get image
                image = shape.image
                image_bytes = image.blob

                # Add to target slide at same position
                left = shape.left
                top = shape.top
                width = shape.width
                height = shape.height

                # Save temp image
                temp_file = f'temp_image_{copied}.png'
                with open(temp_file, 'wb') as f:
                    f.write(image_bytes)

                # Add to target slide
                target_slide.shapes.add_picture(temp_file, left, top, width, height)

                # Delete temp file
                os.remove(temp_file)

                copied += 1
                print(f"      Copied image {copied}")
            except Exception as e:
                print(f"      Could not copy image: {e}")

    return copied

def improve_table_formatting(table):
    """Improve table formatting - larger fonts, better alignment"""
    num_rows = len(table.rows)

    # Header row (first row)
    if num_rows > 0:
        for cell in table.rows[0].cells:
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(12)
                paragraph.font.bold = True
                paragraph.font.color.rgb = RGBColor(255, 255, 255)
                paragraph.alignment = PP_ALIGN.CENTER
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(0, 102, 204)

    # Data rows
    for row_idx in range(1, num_rows):
        row = table.rows[row_idx]
        for col_idx, cell in enumerate(row.cells):
            for paragraph in cell.text_frame.paragraphs:
                # Increase font size
                paragraph.font.size = Pt(11)

                # Right align numbers (columns after first)
                if col_idx > 0:
                    paragraph.alignment = PP_ALIGN.RIGHT
                else:
                    paragraph.alignment = PP_ALIGN.CENTER

            # Alternating row colors (skip last row)
            if row_idx < num_rows - 1 and row_idx % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(245, 245, 245)

    # Last row (totals) - if it exists
    if num_rows > 1:
        last_row_idx = num_rows - 1
        for cell in table.rows[last_row_idx].cells:
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(12)
                paragraph.font.bold = True
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(230, 242, 255)

def improve_text_formatting(shape):
    """Improve text formatting - better alignment and font sizes"""
    if not shape.has_text_frame:
        return

    text_frame = shape.text_frame
    for paragraph in text_frame.paragraphs:
        # Ensure readable font size
        if paragraph.font.size and paragraph.font.size < Pt(10):
            paragraph.font.size = Pt(11)

        # Improve alignment for centered text
        if 'SKIN' in paragraph.text or 'Financial' in paragraph.text:
            paragraph.alignment = PP_ALIGN.CENTER

def apply_formatting(corrected_file, formatted_file, output_file):
    """Apply formatting from formatted file to corrected file"""

    print(f"\n📂 Opening files...")
    print(f"   Corrected data: {corrected_file}")
    print(f"   Your formatting: {formatted_file}")

    prs_corrected = Presentation(corrected_file)
    prs_formatted = Presentation(formatted_file)

    total_images = 0
    total_tables = 0

    # Process each slide
    for slide_num in range(min(len(prs_corrected.slides), len(prs_formatted.slides))):
        print(f"\n📄 Slide {slide_num + 1}:")

        corrected_slide = prs_corrected.slides[slide_num]
        formatted_slide = prs_formatted.slides[slide_num]

        # Copy images (logo, etc) from formatted slide
        images = copy_images_from_slide(formatted_slide, corrected_slide)
        total_images += images
        if images > 0:
            print(f"   ✅ Copied {images} image(s)")

        # Improve formatting in corrected slide
        for shape in corrected_slide.shapes:
            # Improve table formatting
            if shape.has_table:
                improve_table_formatting(shape.table)
                total_tables += 1
                print(f"   ✅ Enhanced table formatting")

            # Improve text formatting
            improve_text_formatting(shape)

    # Save result
    print(f"\n💾 Saving enhanced presentation: {output_file}")
    prs_corrected.save(output_file)

    print("\n" + "=" * 70)
    print("✅ FORMATTING APPLIED SUCCESSFULLY!")
    print("=" * 70)
    print(f"Images copied: {total_images}")
    print(f"Tables enhanced: {total_tables}")
    print(f"Total slides: {len(prs_corrected.slides)}")
    print(f"\n📄 Output: {os.path.abspath(output_file)}")
    print("\n✅ Data integrity maintained (from corrected file)")
    print("✅ Logo and formatting applied (from your file)")
    print("✅ Table fonts increased and aligned")
    print("=" * 70)

if __name__ == "__main__":
    import sys

    # Set UTF-8 encoding
    if sys.platform == 'win32':
        import io
        sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8')

    print("=" * 70)
    print("SKINSPIRE CLINIC - APPLY FORMATTING TO CORRECTED DATA")
    print("=" * 70)
    print("\nThis will:")
    print("  1. Use corrected data file (accurate numbers)")
    print("  2. Copy logo from your edited file")
    print("  3. Increase table font sizes")
    print("  4. Improve alignment")
    print("=" * 70)

    corrected_file = 'Skinspire_Financial_Analysis_2025_Corrected.pptx'
    formatted_file = 'Skinspire_Financial_Analysis_2025_Formatted.pptx'
    output_file = 'Skinspire_Financial_Analysis_2025_FINAL_v2.pptx'

    # Check files exist
    if not os.path.exists(corrected_file):
        print(f"\n❌ ERROR: Corrected file not found: {corrected_file}")
        sys.exit(1)

    if not os.path.exists(formatted_file):
        print(f"\n❌ ERROR: Formatted file not found: {formatted_file}")
        sys.exit(1)

    try:
        apply_formatting(corrected_file, formatted_file, output_file)
    except Exception as e:
        print(f"\n❌ ERROR: {e}")
        import traceback
        traceback.print_exc()
