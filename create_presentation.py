"""
Skinspire HMS - Executive Presentation Generator
Creates a comprehensive PowerPoint presentation for Product Manager review
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from datetime import datetime
import os

# Logo path
LOGO_PATH = r"C:\Users\vinod\AppData\Local\Programs\Skinspire Repository\Skinspire_v2\Tulip.png"

def create_skinspire_presentation():
    # Create presentation with widescreen format
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Color scheme (using RGBColor properly)
    PRIMARY_COLOR = RGBColor(0, 102, 153)       # Teal Blue
    SECONDARY_COLOR = RGBColor(0, 153, 153)     # Cyan
    ACCENT_COLOR = RGBColor(255, 153, 0)        # Orange
    DARK_COLOR = RGBColor(51, 51, 51)           # Dark Gray
    LIGHT_COLOR = RGBColor(240, 240, 240)       # Light Gray
    SUCCESS_COLOR = RGBColor(34, 139, 34)       # Green
    WHITE = RGBColor(255, 255, 255)
    LIGHT_CYAN = RGBColor(200, 230, 240)
    PALE_GREEN = RGBColor(220, 255, 220)
    PALE_YELLOW = RGBColor(255, 250, 205)
    PALE_CYAN = RGBColor(230, 250, 255)

    def add_logo_to_slide(slide):
        """Add Tulip logo to top right corner of the slide"""
        if os.path.exists(LOGO_PATH):
            # Position: top right corner with some padding
            logo_width = Inches(1.0)
            logo_height = Inches(0.8)
            left = prs.slide_width - logo_width - Inches(0.3)
            top = Inches(0.15)
            slide.shapes.add_picture(LOGO_PATH, left, top, logo_width, logo_height)

    def add_title_slide(title, subtitle):
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank

        # Add background shape
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = PRIMARY_COLOR
        shape.line.fill.background()

        # Add title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(12.3), Inches(1.5))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(54)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER

        # Add subtitle
        sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(12.3), Inches(1))
        tf = sub_box.text_frame
        p = tf.paragraphs[0]
        p.text = subtitle
        p.font.size = Pt(28)
        p.font.color.rgb = LIGHT_CYAN
        p.alignment = PP_ALIGN.CENTER

        add_logo_to_slide(slide)
        return slide

    def add_section_slide(title):
        slide = prs.slides.add_slide(prs.slide_layouts[6])

        # Section header background
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(2.5), prs.slide_width, Inches(2.5))
        shape.fill.solid()
        shape.fill.fore_color.rgb = SECONDARY_COLOR
        shape.line.fill.background()

        # Title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(3), Inches(12.3), Inches(1.5))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(48)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER

        add_logo_to_slide(slide)
        return slide

    def add_content_slide(title, content_list, two_columns=False):
        slide = prs.slides.add_slide(prs.slide_layouts[6])

        # Header bar
        header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.3))
        header.fill.solid()
        header.fill.fore_color.rgb = PRIMARY_COLOR
        header.line.fill.background()

        # Title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.3), Inches(0.8))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(32)
        p.font.bold = True
        p.font.color.rgb = WHITE

        if two_columns and len(content_list) >= 2:
            # Left column
            left_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.6), Inches(5.8), Inches(5.5))
            tf = left_box.text_frame
            tf.word_wrap = True
            for i, item in enumerate(content_list[:len(content_list)//2]):
                if i == 0:
                    p = tf.paragraphs[0]
                else:
                    p = tf.add_paragraph()
                p.text = f"• {item}"
                p.font.size = Pt(18)
                p.font.color.rgb = DARK_COLOR
                p.space_after = Pt(12)

            # Right column
            right_box = slide.shapes.add_textbox(Inches(6.8), Inches(1.6), Inches(5.8), Inches(5.5))
            tf = right_box.text_frame
            tf.word_wrap = True
            for i, item in enumerate(content_list[len(content_list)//2:]):
                if i == 0:
                    p = tf.paragraphs[0]
                else:
                    p = tf.add_paragraph()
                p.text = f"• {item}"
                p.font.size = Pt(18)
                p.font.color.rgb = DARK_COLOR
                p.space_after = Pt(12)
        else:
            # Single column
            content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.6), Inches(12.3), Inches(5.5))
            tf = content_box.text_frame
            tf.word_wrap = True
            for i, item in enumerate(content_list):
                if i == 0:
                    p = tf.paragraphs[0]
                else:
                    p = tf.add_paragraph()
                if isinstance(item, tuple):
                    p.text = f"• {item[0]}"
                    p.font.bold = item[1] if len(item) > 1 else False
                else:
                    p.text = f"• {item}"
                p.font.size = Pt(20)
                p.font.color.rgb = DARK_COLOR
                p.space_after = Pt(14)

        add_logo_to_slide(slide)
        return slide

    def add_modules_grid_slide(title, modules):
        slide = prs.slides.add_slide(prs.slide_layouts[6])

        # Header
        header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.3))
        header.fill.solid()
        header.fill.fore_color.rgb = PRIMARY_COLOR
        header.line.fill.background()

        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.3), Inches(0.8))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(32)
        p.font.bold = True
        p.font.color.rgb = WHITE

        # Grid of module boxes
        cols = 4
        box_width = Inches(2.9)
        box_height = Inches(1.3)
        start_x = Inches(0.5)
        start_y = Inches(1.6)
        gap_x = Inches(0.2)
        gap_y = Inches(0.15)

        for i, (module_name, status) in enumerate(modules):
            row = i // cols
            col = i % cols
            x = start_x + col * (box_width + gap_x)
            y = start_y + row * (box_height + gap_y)

            # Module box
            box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, box_width, box_height)
            box.fill.solid()
            if status == "Complete":
                box.fill.fore_color.rgb = PALE_GREEN
            elif status == "In Progress":
                box.fill.fore_color.rgb = PALE_YELLOW
            else:
                box.fill.fore_color.rgb = LIGHT_COLOR
            box.line.color.rgb = PRIMARY_COLOR
            box.line.width = Pt(1.5)

            # Module name
            text_box = slide.shapes.add_textbox(x + Inches(0.1), y + Inches(0.2), box_width - Inches(0.2), Inches(0.6))
            tf = text_box.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = module_name
            p.font.size = Pt(14)
            p.font.bold = True
            p.font.color.rgb = DARK_COLOR
            p.alignment = PP_ALIGN.CENTER

            # Status indicator
            status_box = slide.shapes.add_textbox(x + Inches(0.1), y + Inches(0.85), box_width - Inches(0.2), Inches(0.3))
            tf = status_box.text_frame
            p = tf.paragraphs[0]
            p.text = status
            p.font.size = Pt(11)
            if status == "Complete":
                p.font.color.rgb = SUCCESS_COLOR
            elif status == "In Progress":
                p.font.color.rgb = ACCENT_COLOR
            else:
                p.font.color.rgb = DARK_COLOR
            p.alignment = PP_ALIGN.CENTER

        add_logo_to_slide(slide)
        return slide

    def add_architecture_slide():
        slide = prs.slides.add_slide(prs.slide_layouts[6])

        # Header
        header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.3))
        header.fill.solid()
        header.fill.fore_color.rgb = PRIMARY_COLOR
        header.line.fill.background()

        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.3), Inches(0.8))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = "System Architecture"
        p.font.size = Pt(32)
        p.font.bold = True
        p.font.color.rgb = WHITE

        # Architecture layers
        layers = [
            ("Presentation Layer", "HTML/CSS/JavaScript, Jinja2 Templates, TailwindCSS", Inches(1.5)),
            ("Application Layer", "Flask Framework, Views, API Routes, Services", Inches(2.7)),
            ("Business Logic", "Universal Engine, RBAC, Discount Service, Billing Service", Inches(3.9)),
            ("Data Access", "SQLAlchemy ORM, Database Service, Repositories", Inches(5.1)),
            ("Database", "PostgreSQL with JSONB support, Redis Cache", Inches(6.3)),
        ]

        for name, desc, y in layers:
            # Layer box
            box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.5), y, Inches(10.3), Inches(1))
            box.fill.solid()
            box.fill.fore_color.rgb = SECONDARY_COLOR
            box.line.color.rgb = PRIMARY_COLOR
            box.line.width = Pt(2)

            # Layer name
            name_box = slide.shapes.add_textbox(Inches(1.7), y + Inches(0.15), Inches(3), Inches(0.4))
            tf = name_box.text_frame
            p = tf.paragraphs[0]
            p.text = name
            p.font.size = Pt(16)
            p.font.bold = True
            p.font.color.rgb = WHITE

            # Description
            desc_box = slide.shapes.add_textbox(Inches(4.8), y + Inches(0.15), Inches(6.8), Inches(0.7))
            tf = desc_box.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = desc
            p.font.size = Pt(13)
            p.font.color.rgb = PALE_CYAN

        add_logo_to_slide(slide)
        return slide

    def add_multi_tenant_slide():
        slide = prs.slides.add_slide(prs.slide_layouts[6])

        # Header
        header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.3))
        header.fill.solid()
        header.fill.fore_color.rgb = PRIMARY_COLOR
        header.line.fill.background()

        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.3), Inches(0.8))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = "Multi-Tenant Architecture"
        p.font.size = Pt(32)
        p.font.bold = True
        p.font.color.rgb = WHITE

        # Hospital Level Box
        hospital_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(1.5), Inches(12.3), Inches(1.2))
        hospital_box.fill.solid()
        hospital_box.fill.fore_color.rgb = PRIMARY_COLOR
        hospital_box.line.color.rgb = RGBColor(0, 80, 120)
        hospital_box.line.width = Pt(3)

        h_text = slide.shapes.add_textbox(Inches(0.7), Inches(1.7), Inches(11.9), Inches(0.9))
        tf = h_text.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = "HOSPITAL LEVEL (Tenant Isolation)"
        p.font.size = Pt(20)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p = tf.add_paragraph()
        p.text = "Complete data isolation | Separate configurations | Independent GL accounts | Custom branding"
        p.font.size = Pt(14)
        p.font.color.rgb = LIGHT_CYAN

        # Branch containers (3 branches)
        branch_colors = [RGBColor(0, 130, 130), RGBColor(0, 150, 150), RGBColor(0, 170, 170)]
        branch_names = ["Branch 1: Main Clinic", "Branch 2: City Center", "Branch 3: Mall Outlet"]
        branch_data = [
            "Users: 15 | Patients: 2,500",
            "Users: 8 | Patients: 1,200",
            "Users: 5 | Patients: 800"
        ]

        for i in range(3):
            x = Inches(0.7 + i * 4.1)
            branch = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(2.9), Inches(3.9), Inches(2.2))
            branch.fill.solid()
            branch.fill.fore_color.rgb = branch_colors[i]
            branch.line.color.rgb = PRIMARY_COLOR
            branch.line.width = Pt(2)

            b_text = slide.shapes.add_textbox(x + Inches(0.1), Inches(3.0), Inches(3.7), Inches(2.0))
            tf = b_text.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = branch_names[i]
            p.font.size = Pt(16)
            p.font.bold = True
            p.font.color.rgb = WHITE
            p = tf.add_paragraph()
            p.text = branch_data[i]
            p.font.size = Pt(12)
            p.font.color.rgb = PALE_CYAN
            p = tf.add_paragraph()
            p.text = "Branch-specific: Inventory, Staff, Appointments"
            p.font.size = Pt(11)
            p.font.color.rgb = PALE_CYAN

        # Key Features
        features_box = slide.shapes.add_textbox(Inches(0.5), Inches(5.3), Inches(12.3), Inches(2.0))
        tf = features_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = "Key Capabilities:"
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = DARK_COLOR

        features = [
            "Hospital-level: Shared masters (Services, Medicines, Suppliers), Consolidated reports, Central configuration",
            "Branch-level: Independent inventory, Local staff assignments, Branch-specific promotions, Separate cash registers",
            "Cross-branch: Patient records accessible across branches, Inter-branch transfers, Unified loyalty program"
        ]
        for feat in features:
            p = tf.add_paragraph()
            p.text = f"• {feat}"
            p.font.size = Pt(13)
            p.font.color.rgb = DARK_COLOR

        add_logo_to_slide(slide)
        return slide

    def add_integration_diagram_slide():
        slide = prs.slides.add_slide(prs.slide_layouts[6])

        # Header
        header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.3))
        header.fill.solid()
        header.fill.fore_color.rgb = PRIMARY_COLOR
        header.line.fill.background()

        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.3), Inches(0.8))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = "External System Integrations"
        p.font.size = Pt(32)
        p.font.bold = True
        p.font.color.rgb = WHITE

        # Center - SkinSpire HMS
        center_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(5), Inches(3), Inches(3.3), Inches(1.5))
        center_box.fill.solid()
        center_box.fill.fore_color.rgb = PRIMARY_COLOR
        center_box.line.color.rgb = RGBColor(0, 80, 120)
        center_box.line.width = Pt(3)

        c_text = slide.shapes.add_textbox(Inches(5.1), Inches(3.3), Inches(3.1), Inches(1.0))
        tf = c_text.text_frame
        p = tf.paragraphs[0]
        p.text = "SkinSpire HMS"
        p.font.size = Pt(20)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER
        p = tf.add_paragraph()
        p.text = "Core Platform"
        p.font.size = Pt(14)
        p.font.color.rgb = LIGHT_CYAN
        p.alignment = PP_ALIGN.CENTER

        # Integration boxes around center
        integrations = [
            ("Paytm EDC", "Payment Gateway\nCard, UPI, QR", Inches(1), Inches(1.5), RGBColor(0, 175, 240)),
            ("ICICI Bank", "Corporate Banking\nAuto Reconciliation", Inches(9.8), Inches(1.5), RGBColor(255, 102, 0)),
            ("Material Scanner", "Barcode/QR\nInventory Mgmt", Inches(1), Inches(4.8), RGBColor(76, 175, 80)),
            ("Voice Assistant", "AI Documentation\nHands-free Input", Inches(9.8), Inches(4.8), RGBColor(156, 39, 176)),
        ]

        for name, desc, x, y, color in integrations:
            box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(2.8), Inches(1.3))
            box.fill.solid()
            box.fill.fore_color.rgb = color
            box.line.color.rgb = DARK_COLOR
            box.line.width = Pt(1.5)

            i_text = slide.shapes.add_textbox(x + Inches(0.1), y + Inches(0.15), Inches(2.6), Inches(1.0))
            tf = i_text.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = name
            p.font.size = Pt(15)
            p.font.bold = True
            p.font.color.rgb = WHITE
            p.alignment = PP_ALIGN.CENTER
            p = tf.add_paragraph()
            p.text = desc
            p.font.size = Pt(11)
            p.font.color.rgb = WHITE
            p.alignment = PP_ALIGN.CENTER

        # Connection lines description
        desc_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.3), Inches(12.3), Inches(1.0))
        tf = desc_box.text_frame
        p = tf.paragraphs[0]
        p.text = "All integrations use secure REST APIs with OAuth 2.0 authentication and end-to-end encryption"
        p.font.size = Pt(14)
        p.font.color.rgb = DARK_COLOR
        p.alignment = PP_ALIGN.CENTER

        add_logo_to_slide(slide)
        return slide

    def add_ai_features_slide():
        slide = prs.slides.add_slide(prs.slide_layouts[6])

        # Header
        header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.3))
        header.fill.solid()
        header.fill.fore_color.rgb = RGBColor(128, 0, 128)  # Purple for AI
        header.line.fill.background()

        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.3), Inches(0.8))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = "AI-Powered Features Roadmap"
        p.font.size = Pt(32)
        p.font.bold = True
        p.font.color.rgb = WHITE

        # AI Feature boxes
        ai_features = [
            ("Voice Assistant", "Clinical Documentation",
             ["Record patient history hands-free", "Voice-to-prescription conversion", "Automated consultation notes", "Multi-language support (EN/HI)"],
             Inches(0.5), RGBColor(156, 39, 176)),
            ("Smart Reconciliation", "Account Management",
             ["Auto-match bank transactions", "Identify payment discrepancies", "Suggest GL postings", "Anomaly detection alerts"],
             Inches(4.5), RGBColor(63, 81, 181)),
            ("Analytics Engine", "Business Intelligence",
             ["Promotion effectiveness scoring", "Patient churn prediction", "Inventory demand forecasting", "Revenue optimization insights"],
             Inches(8.5), RGBColor(0, 150, 136)),
        ]

        for title, subtitle, items, x, color in ai_features:
            box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(1.5), Inches(3.8), Inches(4.2))
            box.fill.solid()
            box.fill.fore_color.rgb = color
            box.line.color.rgb = DARK_COLOR
            box.line.width = Pt(2)

            f_text = slide.shapes.add_textbox(x + Inches(0.15), Inches(1.6), Inches(3.5), Inches(4.0))
            tf = f_text.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = title
            p.font.size = Pt(18)
            p.font.bold = True
            p.font.color.rgb = WHITE
            p = tf.add_paragraph()
            p.text = subtitle
            p.font.size = Pt(13)
            p.font.color.rgb = PALE_CYAN
            p.space_after = Pt(8)

            for item in items:
                p = tf.add_paragraph()
                p.text = f"• {item}"
                p.font.size = Pt(12)
                p.font.color.rgb = WHITE

        # Timeline
        timeline_box = slide.shapes.add_textbox(Inches(0.5), Inches(5.9), Inches(12.3), Inches(1.3))
        tf = timeline_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = "Implementation Timeline"
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = DARK_COLOR
        p = tf.add_paragraph()
        p.text = "Q1 2026: Voice Assistant Beta | Q2 2026: Smart Reconciliation | Q3 2026: Predictive Analytics | Q4 2026: Full AI Suite"
        p.font.size = Pt(13)
        p.font.color.rgb = DARK_COLOR

        add_logo_to_slide(slide)
        return slide

    def add_timeline_slide():
        slide = prs.slides.add_slide(prs.slide_layouts[6])

        # Header
        header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.3))
        header.fill.solid()
        header.fill.fore_color.rgb = PRIMARY_COLOR
        header.line.fill.background()

        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.3), Inches(0.8))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = "Project Timeline"
        p.font.size = Pt(32)
        p.font.bold = True
        p.font.color.rgb = WHITE

        # Timeline milestones
        milestones = [
            ("Feb 2025", "Planning & Design", "Requirements, Architecture, UI/UX Design", RGBColor(63, 81, 181)),
            ("May 2025", "Development Begin", "Core modules, Database, Security framework", RGBColor(0, 150, 136)),
            ("Jan 2026", "Phase 1: Core", "Patient, Billing, Inventory, Payments go-live", RGBColor(76, 175, 80)),
            ("Apr 2026", "Phase 2: AI & Integration", "Voice Assistant, Bank/Payment Gateway", RGBColor(255, 152, 0)),
            ("Jul 2026", "Phase 3: Pilot", "Pilot deployment @ SkinSpire clinic", RGBColor(233, 30, 99)),
            ("Sep 2026", "Phase 4: Roll Out", "Multi-tenant hosting for 5 clinics", RGBColor(156, 39, 176)),
        ]

        # Draw timeline line
        line_y = Inches(4.2)
        line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), line_y, Inches(11.7), Inches(0.08))
        line.fill.solid()
        line.fill.fore_color.rgb = DARK_COLOR
        line.line.fill.background()

        # Draw milestones
        start_x = Inches(0.8)
        spacing = Inches(2.1)

        for i, (date, title, desc, color) in enumerate(milestones):
            x = start_x + (i * spacing)

            # Circle marker on timeline
            circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(0.85), line_y - Inches(0.15), Inches(0.38), Inches(0.38))
            circle.fill.solid()
            circle.fill.fore_color.rgb = color
            circle.line.color.rgb = WHITE
            circle.line.width = Pt(2)

            # Alternating above/below timeline
            if i % 2 == 0:
                # Above timeline
                box_y = Inches(1.6)
                connector_start_y = Inches(3.3)
                connector_end_y = line_y
            else:
                # Below timeline
                box_y = Inches(4.7)
                connector_start_y = line_y + Inches(0.1)
                connector_end_y = Inches(4.65)

            # Milestone box
            box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, box_y, Inches(2.0), Inches(1.6))
            box.fill.solid()
            box.fill.fore_color.rgb = color
            box.line.color.rgb = DARK_COLOR
            box.line.width = Pt(1)

            # Date
            date_box = slide.shapes.add_textbox(x + Inches(0.08), box_y + Inches(0.08), Inches(1.84), Inches(0.35))
            tf = date_box.text_frame
            p = tf.paragraphs[0]
            p.text = date
            p.font.size = Pt(13)
            p.font.bold = True
            p.font.color.rgb = WHITE
            p.alignment = PP_ALIGN.CENTER

            # Title
            title_box = slide.shapes.add_textbox(x + Inches(0.08), box_y + Inches(0.4), Inches(1.84), Inches(0.5))
            tf = title_box.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = title
            p.font.size = Pt(12)
            p.font.bold = True
            p.font.color.rgb = WHITE
            p.alignment = PP_ALIGN.CENTER

            # Description
            desc_box = slide.shapes.add_textbox(x + Inches(0.08), box_y + Inches(0.9), Inches(1.84), Inches(0.65))
            tf = desc_box.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = desc
            p.font.size = Pt(9)
            p.font.color.rgb = RGBColor(240, 240, 240)
            p.alignment = PP_ALIGN.CENTER

        # Legend / Note
        note_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.5), Inches(12.3), Inches(0.7))
        tf = note_box.text_frame
        p = tf.paragraphs[0]
        p.text = "Timeline: 20 months from planning to multi-clinic deployment | Current Status: Development Phase"
        p.font.size = Pt(14)
        p.font.color.rgb = DARK_COLOR
        p.alignment = PP_ALIGN.CENTER

        add_logo_to_slide(slide)
        return slide

    def add_app_workflow_slide():
        slide = prs.slides.add_slide(prs.slide_layouts[6])

        # Header
        header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.3))
        header.fill.solid()
        header.fill.fore_color.rgb = PRIMARY_COLOR
        header.line.fill.background()

        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.3), Inches(0.8))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = "Multi-App Ecosystem - Self Service & Workflow"
        p.font.size = Pt(30)
        p.font.bold = True
        p.font.color.rgb = WHITE

        # Three App columns
        apps = [
            ("Patient App", RGBColor(76, 175, 80), [
                "Self-Service Appointments",
                "Book consultations online",
                "Schedule procedures/treatments",
                "View & reschedule bookings",
                "",
                "Health Records (EMR)",
                "Access medical history",
                "View prescriptions & reports",
                "Download invoices & receipts",
                "",
                "Payments & Loyalty",
                "Pay bills online",
                "Track loyalty points",
                "View promotion offers"
            ]),
            ("Doctor App", RGBColor(33, 150, 243), [
                "Clinical Workflow",
                "View daily appointments",
                "Patient queue management",
                "Access patient EMR history",
                "",
                "Consultation & EMR",
                "Record vitals & complaints",
                "Voice-enabled note taking",
                "Digital prescriptions",
                "",
                "Approvals & Reviews",
                "Approve treatment plans",
                "Review lab results",
                "Authorize procedures"
            ]),
            ("Staff App", RGBColor(255, 152, 0), [
                "Front Desk Operations",
                "Check-in patients",
                "Manage walk-ins",
                "Payment collection",
                "",
                "Inventory & Billing",
                "Stock verification",
                "Generate invoices",
                "Apply discounts/promos",
                "",
                "Workflow & Approvals",
                "Request manager approvals",
                "Process refunds/credits",
                "Handle escalations"
            ])
        ]

        for i, (app_name, color, features) in enumerate(apps):
            x = Inches(0.5 + i * 4.2)

            # App header box
            app_header = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(1.5), Inches(4.0), Inches(0.7))
            app_header.fill.solid()
            app_header.fill.fore_color.rgb = color
            app_header.line.fill.background()

            name_box = slide.shapes.add_textbox(x + Inches(0.1), Inches(1.6), Inches(3.8), Inches(0.5))
            tf = name_box.text_frame
            p = tf.paragraphs[0]
            p.text = app_name
            p.font.size = Pt(18)
            p.font.bold = True
            p.font.color.rgb = WHITE
            p.alignment = PP_ALIGN.CENTER

            # Features box
            feat_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(2.25), Inches(4.0), Inches(4.6))
            feat_box.fill.solid()
            feat_box.fill.fore_color.rgb = RGBColor(250, 250, 250)
            feat_box.line.color.rgb = color
            feat_box.line.width = Pt(2)

            # Features text
            text_box = slide.shapes.add_textbox(x + Inches(0.15), Inches(2.35), Inches(3.7), Inches(4.4))
            tf = text_box.text_frame
            tf.word_wrap = True

            first_para = True
            for feat in features:
                if first_para:
                    p = tf.paragraphs[0]
                    first_para = False
                else:
                    p = tf.add_paragraph()

                if feat == "":
                    p.text = ""
                    p.space_after = Pt(4)
                elif feat in ["Self-Service Appointments", "Health Records (EMR)", "Payments & Loyalty",
                             "Clinical Workflow", "Consultation & EMR", "Approvals & Reviews",
                             "Front Desk Operations", "Inventory & Billing", "Workflow & Approvals"]:
                    p.text = feat
                    p.font.size = Pt(11)
                    p.font.bold = True
                    p.font.color.rgb = color
                    p.space_before = Pt(4)
                else:
                    p.text = f"• {feat}"
                    p.font.size = Pt(10)
                    p.font.color.rgb = DARK_COLOR

        add_logo_to_slide(slide)
        return slide

    def add_business_insights_slide():
        slide = prs.slides.add_slide(prs.slide_layouts[6])

        # Header
        header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.3))
        header.fill.solid()
        header.fill.fore_color.rgb = RGBColor(63, 81, 181)  # Indigo for analytics
        header.line.fill.background()

        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.3), Inches(0.8))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = "Business Insights & Analytics Dashboard"
        p.font.size = Pt(30)
        p.font.bold = True
        p.font.color.rgb = WHITE

        # Analytics categories in 2x3 grid
        insights = [
            ("P&L Analysis", "Financial Performance",
             ["Revenue by service category", "Cost center tracking", "Gross margin analysis", "Month-over-month trends"],
             RGBColor(76, 175, 80), Inches(0.5), Inches(1.5)),

            ("Service Profitability", "Service-wise Analysis",
             ["Revenue per service type", "Cost per procedure", "Therapist productivity", "Package vs individual sales"],
             RGBColor(0, 150, 136), Inches(4.5), Inches(1.5)),

            ("Cash Flow Planning", "Treasury Management",
             ["Daily cash position", "Receivables aging", "Payment forecasting", "Branch-wise cash flow"],
             RGBColor(33, 150, 243), Inches(8.5), Inches(1.5)),

            ("Inventory Analytics", "Stock Intelligence",
             ["Stock turnover ratio", "Expiry alerts & wastage", "Reorder optimization", "Supplier performance"],
             RGBColor(255, 152, 0), Inches(0.5), Inches(4.2)),

            ("Patient Analytics", "Customer Insights",
             ["Patient lifetime value", "Visit frequency analysis", "Treatment preferences", "Churn risk scoring"],
             RGBColor(233, 30, 99), Inches(4.5), Inches(4.2)),

            ("Campaign & Product", "Marketing ROI",
             ["Campaign effectiveness", "Discount impact analysis", "Product profitability", "Promotion conversion rates"],
             RGBColor(156, 39, 176), Inches(8.5), Inches(4.2)),
        ]

        for title, subtitle, items, color, x, y in insights:
            # Card box
            card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(3.8), Inches(2.5))
            card.fill.solid()
            card.fill.fore_color.rgb = color
            card.line.fill.background()

            # Title
            t_box = slide.shapes.add_textbox(x + Inches(0.15), y + Inches(0.1), Inches(3.5), Inches(0.4))
            tf = t_box.text_frame
            p = tf.paragraphs[0]
            p.text = title
            p.font.size = Pt(15)
            p.font.bold = True
            p.font.color.rgb = WHITE

            # Subtitle
            s_box = slide.shapes.add_textbox(x + Inches(0.15), y + Inches(0.45), Inches(3.5), Inches(0.3))
            tf = s_box.text_frame
            p = tf.paragraphs[0]
            p.text = subtitle
            p.font.size = Pt(11)
            p.font.color.rgb = RGBColor(220, 220, 220)

            # Items
            i_box = slide.shapes.add_textbox(x + Inches(0.15), y + Inches(0.8), Inches(3.5), Inches(1.6))
            tf = i_box.text_frame
            tf.word_wrap = True
            for idx, item in enumerate(items):
                if idx == 0:
                    p = tf.paragraphs[0]
                else:
                    p = tf.add_paragraph()
                p.text = f"• {item}"
                p.font.size = Pt(11)
                p.font.color.rgb = WHITE

        # Footer note
        note_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.85), Inches(12.3), Inches(0.5))
        tf = note_box.text_frame
        p = tf.paragraphs[0]
        p.text = "Real-time dashboards with drill-down capabilities | Export to Excel/PDF | Scheduled email reports"
        p.font.size = Pt(13)
        p.font.color.rgb = DARK_COLOR
        p.alignment = PP_ALIGN.CENTER

        add_logo_to_slide(slide)
        return slide

    # ==================== CREATE SLIDES ====================

    # Slide 1: Title
    add_title_slide(
        "SkinSpire Clinic HMS",
        "Healthcare Management System - Executive Overview\nNovember 2025"
    )

    # Slide 2: Executive Summary
    add_content_slide("Executive Summary", [
        "SkinSpire is a comprehensive Hospital Management System (HMS) designed for mid-size dermatology and aesthetic clinics",
        "Built with enterprise-grade architecture supporting multi-hospital, multi-branch operations",
        "Fully compliant with Indian GST regulations and healthcare industry standards",
        "Features advanced Role-Based Access Control (RBAC) with branch-level security",
        "Innovative 'Universal Engine' reduces new module development time by 90%",
        "Production-ready with comprehensive billing, inventory, and financial management"
    ])

    # Slide 3: Vision
    add_section_slide("Our Vision")

    # Slide 4: Vision Statement
    add_content_slide("Vision - Transforming Healthcare Operations", [
        ("\"To become the leading healthcare management platform for dermatology and aesthetic clinics across India\"", True),
        "Empower clinics with intelligent automation to focus on patient care, not paperwork",
        "Create a unified ecosystem connecting patients, providers, and financial systems",
        "Enable data-driven decision making through AI-powered analytics and insights",
        "Build a scalable platform that grows from single clinic to multi-city hospital chains",
        "Pioneer voice-enabled healthcare documentation for hands-free clinical workflows",
        "Achieve seamless integration with India's digital payment and banking infrastructure"
    ])

    # Slide 5: Strategic Goals
    add_content_slide("Strategic Goals 2025-2027", [
        "Year 1: Complete core HMS with payment gateway and bank integrations",
        "Year 1: Deploy AI-powered voice assistant for clinical documentation",
        "Year 2: Launch patient mobile app with appointment and payment features",
        "Year 2: Implement predictive analytics for inventory and promotions",
        "Year 3: Expand to 100+ clinic deployments across major Indian cities",
        "Year 3: Achieve EMR/EHR certification and insurance integration",
        "Ongoing: Maintain 99.9% uptime with enterprise-grade security"
    ])

    # Slide 6: Project Timeline
    add_timeline_slide()

    # Slide 7: Business Value Proposition
    add_section_slide("Business Value Proposition")

    # Slide 4: Key Business Benefits
    add_content_slide("Key Business Benefits", [
        "Multi-Tenant Architecture: Single deployment serves multiple hospitals/branches",
        "Regulatory Compliance: Full GST support (CGST, SGST, IGST) for Indian healthcare",
        "Financial Integration: Double-entry accounting with Chart of Accounts",
        "Scalable Design: From single clinic to hospital chain without code changes",
        "Security First: Field-level encryption, audit trails, and RBAC",
        "Cost Efficient: Reduces operational overhead with streamlined workflows",
        "Future Ready: Foundation for EMR/EHR, telemedicine, and insurance integration"
    ])

    # Slide 5: Core Modules Overview
    add_section_slide("Core Modules")

    # Slide 6: Module Status Grid
    modules = [
        ("Patient Management", "Complete"),
        ("Appointment Booking", "Complete"),
        ("Doctor Consultation", "Complete"),
        ("Billing & Invoicing", "Complete"),
        ("Payment Processing", "Complete"),
        ("Advance Payments", "Complete"),
        ("Purchase Orders", "Complete"),
        ("Supplier Invoices", "Complete"),
        ("Supplier Payments", "Complete"),
        ("Inventory Management", "Complete"),
        ("Promotions & Discounts", "Complete"),
        ("Universal Engine", "Complete"),
        ("RBAC Security", "Complete"),
        ("GL Accounting", "Complete"),
        ("Credit Notes", "Complete"),
        ("Reports & Analytics", "In Progress"),
    ]
    add_modules_grid_slide("Module Implementation Status", modules)

    # Slide 7: Patient Journey
    add_content_slide("Patient Journey - End to End", [
        "Registration: Complete patient profile with documents and medical history",
        "Appointment: Online/walk-in booking with doctor/service selection",
        "Consultation: Vitals recording, diagnosis, prescriptions, follow-ups",
        "Services: Package selection, treatment sessions, therapist assignments",
        "Billing: Multi-type invoices (Services, Medicines, Misc) with GST",
        "Payments: Cash, Card, UPI, Advance payments, and settlement",
        "Loyalty: Points accumulation, tier progression, promotional discounts"
    ])

    # Slide: Multi-App Ecosystem
    add_app_workflow_slide()

    # Slide: Business Insights & Analytics
    add_business_insights_slide()

    # Slide 8: Billing & Financial
    add_content_slide("Billing & Financial Management", [
        "Invoice Types: Services, Prescription Drugs, OTC/Cosmetics, Miscellaneous",
        "GST Compliance: Automatic CGST/SGST/IGST calculation based on state codes",
        "Multiple Payment Methods: Cash, Credit Card, Debit Card, UPI supported",
        "Advance Payment System: Patient deposits with automatic invoice allocation",
        "Discount Framework: Campaign, Bulk, Loyalty, VIP with configurable stacking",
        "Credit Notes: Full and partial refunds with GL integration",
        "Chart of Accounts: Complete double-entry accounting system"
    ])

    # Slide 9: Promotions & Discounts
    add_content_slide("Promotions & Discounts System", [
        "Campaign Management: Time-bound promotions with approval workflow",
        "Discount Types: Percentage, Fixed Amount, Buy X Get Y Free",
        "Stacking Configuration: Exclusive, Incremental, Absolute modes",
        "Patient Targeting: VIP-specific, loyalty tier-based, personalized offers",
        "Campaign Groups: Organize related campaigns for easy management",
        "Timeline Visualization: Visual campaign planning with holiday awareness",
        "Simulation Tool: Preview discount impact before applying to invoices"
    ])

    # Slide: Technical Architecture Section
    add_section_slide("Technical Architecture")

    # Slide: Architecture Diagram
    add_architecture_slide()

    # Slide: Multi-Tenant Architecture
    add_multi_tenant_slide()

    # Slide: Technology Stack
    add_content_slide("Technology Stack", [
        "Backend: Python 3.13 with Flask 3.1 web framework",
        "Database: PostgreSQL with JSONB for flexible data structures",
        "ORM: SQLAlchemy 2.0 with Alembic migrations",
        "Caching: Redis for session management and performance",
        "Frontend: Jinja2 templates with TailwindCSS styling",
        "Security: Cryptography library, PyJWT for tokens",
        "PDF Generation: ReportLab for invoices and reports",
        "Testing: Pytest with comprehensive test coverage"
    ], two_columns=True)

    # Slide 13: Universal Engine
    add_content_slide("Universal Engine Innovation", [
        "Problem: Traditional development requires 18+ hours per new module",
        "Solution: Configuration-driven architecture with reusable components",
        "Result: New entities can be added in 30 minutes (90% reduction)",
        "Features: Auto-generated list views, search, pagination, export",
        "Benefits: Consistent UI, reduced bugs, easier maintenance",
        "Status: Successfully deployed for Supplier Payments, ready for all entities"
    ])

    # Slide: Security & Compliance
    add_content_slide("Security & Compliance", [
        "Role-Based Access Control: Granular permissions for all operations",
        "Branch-Level Security: Users see only their assigned branch data",
        "Hospital Isolation: Complete data separation in multi-tenant setup",
        "Field Encryption: Sensitive patient data encrypted at rest",
        "Audit Trails: Complete logging of who created/modified records",
        "Session Management: Redis-based secure session handling",
        "CSRF Protection: All forms protected against cross-site attacks"
    ])

    # Slide: External Integrations Section
    add_section_slide("External System Integrations")

    # Slide: Integration Diagram
    add_integration_diagram_slide()

    # Slide: Paytm EDC Integration
    add_content_slide("Paytm EDC Payment Gateway", [
        "Unified Payment Processing: Single device for Cards, UPI, and QR payments",
        "Real-time Transaction Sync: Automatic invoice status updates on payment",
        "Multi-terminal Support: Multiple EDC devices per branch with central tracking",
        "Settlement Reports: Daily reconciliation with transaction-level details",
        "Refund Processing: Seamless refund initiation from credit notes",
        "Offline Mode: Store-and-forward for network interruptions",
        "PCI-DSS Compliant: Secure card data handling without local storage"
    ])

    # Slide: ICICI Bank Integration
    add_content_slide("ICICI Bank Corporate Integration", [
        "Account Statement Sync: Automated daily/weekly bank statement import",
        "Payment Reconciliation: Auto-match incoming payments with invoices",
        "Supplier Payments: Direct bank transfer for vendor settlements",
        "Cash Management: Real-time cash position across all branches",
        "Bulk Payments: Process multiple supplier payments in single batch",
        "MIS Reports: Bank transaction analysis and cash flow reporting",
        "API Security: Bank-grade encryption with certificate authentication"
    ])

    # Slide: Material Scanner Integration
    add_content_slide("Material Scanner & Inventory", [
        "Barcode Scanning: Fast product lookup during billing and receiving",
        "QR Code Support: Track batches and expiry with 2D codes",
        "Stock Taking: Rapid physical inventory verification",
        "Goods Receiving: Scan-based verification of supplier deliveries",
        "Medicine Dispensing: Verify correct items during pharmacy billing",
        "Hardware Support: Compatible with handheld and fixed scanners",
        "Batch Tracking: Full traceability from receipt to sale"
    ])

    # Slide: Voice Assistant Overview
    add_content_slide("Voice Assistant Integration", [
        "Hands-free Documentation: Record consultations without typing",
        "Patient History: Voice capture of complaints, symptoms, allergies",
        "Prescription Dictation: Speak medication orders for auto-conversion",
        "Multi-language: Support for English, Hindi, and regional languages",
        "Context Awareness: Understand medical terminology and drug names",
        "Privacy Controls: On-premise processing option for sensitive data",
        "Integration: Direct flow into consultation notes and prescriptions"
    ])

    # Slide: AI Features Section
    add_section_slide("AI-Powered Future")

    # Slide: AI Features Roadmap
    add_ai_features_slide()

    # Slide: AI Use Cases Detail
    add_content_slide("AI Use Cases in Detail", [
        ("Account Reconciliation AI", True),
        "Automatically match bank credits to patient payments using amount, date, and reference patterns",
        "Flag unusual transactions for manual review with confidence scores",
        ("Voice-to-EMR Conversion", True),
        "Convert doctor's verbal notes into structured EMR entries (complaints, diagnosis, treatment)",
        ("Promotion Effectiveness AI", True),
        "Analyze campaign performance, predict ROI, recommend optimal discount structures",
        ("Inventory Intelligence", True),
        "Predict stock requirements based on historical patterns and upcoming appointments"
    ])

    # Slide: What's Been Achieved
    add_section_slide("Project Achievements")

    # Slide 16: Completed Features
    add_content_slide("Completed Features (100%)", [
        "Core Patient Management with document storage",
        "Complete Billing workflow with GST compliance",
        "Multi-method Payment Processing",
        "Advance Payment and Settlement System",
        "Full Inventory and Purchase Order Management",
        "Supplier Invoice and Payment workflows",
        "Comprehensive Promotions and Discounts module",
        "Universal View Engine for rapid development",
        "RBAC with multi-branch support"
    ], two_columns=True)

    # Slide 17: Quality Metrics
    add_content_slide("Quality & Architecture Metrics", [
        ("Enterprise-Grade Architecture: A+ Rating", True),
        "Separation of Concerns with clear layer boundaries",
        "Configuration-driven behavior minimizing code changes",
        "Comprehensive error handling with graceful fallbacks",
        "Zero code duplication through Universal Engine",
        ("Backward Compatibility: 100%", True),
        "All existing functionality preserved during enhancements",
        ("Production Readiness: Immediate Deployment", True)
    ])

    # Slide 18: Next Steps
    add_section_slide("Roadmap & Next Steps")

    # Slide 19: Recommended Next Steps
    add_content_slide("Recommended Next Steps", [
        "Phase 1: Complete Reports & Analytics dashboard",
        "Phase 2: Implement EDC/Payment Gateway integration (Paytm)",
        "Phase 3: Build RBAC Administration Interface",
        "Phase 4: Add Email/SMS Notification System",
        "Phase 5: Develop Patient Portal for self-service",
        "Phase 6: EMR/EHR Integration foundation",
        "Phase 7: Insurance claim processing module"
    ])

    # Slide 20: Investment Areas
    add_content_slide("Areas Requiring Investment", [
        "User Training: Staff onboarding and documentation",
        "Data Migration: If replacing existing systems",
        "Infrastructure: Production server setup and security hardening",
        "Testing: User Acceptance Testing (UAT) with real scenarios",
        "Support: Helpdesk setup and issue tracking process",
        "Compliance: HIPAA/local healthcare regulation audit",
        "Backup: Disaster recovery and business continuity planning"
    ])

    # Slide 21: Summary
    add_content_slide("Summary", [
        ("SkinSpire HMS is production-ready for mid-size clinics", True),
        "Comprehensive coverage of patient, billing, and inventory workflows",
        "Enterprise-grade security with RBAC and audit compliance",
        "Innovative Universal Engine for rapid future development",
        "Flexible promotions system to drive patient engagement",
        "Strong foundation for growth: multi-branch, multi-hospital ready",
        ("Ready to transform clinic operations!", True)
    ])

    # Slide 22: Thank You
    add_title_slide(
        "Thank You",
        "Questions & Discussion\n\nSkinSpire HMS - Built for Healthcare Excellence"
    )

    # Save presentation
    output_path = "SkinSpire_HMS_Executive_Overview.pptx"
    prs.save(output_path)
    print(f"Presentation saved to: {output_path}")
    return output_path

if __name__ == "__main__":
    create_skinspire_presentation()
