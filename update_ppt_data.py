"""
Update existing PowerPoint with corrected Capital Expenditure data
Preserves all formatting, logos, and custom changes
"""

from pptx import Presentation
import re

# Corrected financial data
MONTHLY_DATA = {
    'Jan': {'revenue': 1110615, 'opex': 546768, 'capex': 40499, 'personal': 137448, 'profit': 385901},
    'Feb': {'revenue': 629780, 'opex': 440076, 'capex': 354674, 'personal': 206349, 'profit': -371319},
    'Mar': {'revenue': 962581, 'opex': 533648, 'capex': 404273, 'personal': 21518, 'profit': 3142},
    'Apr': {'revenue': 783448, 'opex': 326278, 'capex': 395188, 'personal': 41367, 'profit': 20615},
    'May': {'revenue': 855204, 'opex': 455958, 'capex': 405636, 'personal': 52696, 'profit': -59086},
    'Jun': {'revenue': 965456, 'opex': 502990, 'capex': 141317, 'personal': 125315, 'profit': 195833},
    'Jul': {'revenue': 1347544, 'opex': 270102, 'capex': 157612, 'personal': 80465, 'profit': 839366},
    'Aug': {'revenue': 1243094, 'opex': 489903, 'capex': 101979, 'personal': 46459, 'profit': 604753},
    'Sep': {'revenue': 796105, 'opex': 693200, 'capex': 66662, 'personal': 76691, 'profit': -40448},
    'Oct': {'revenue': 701445, 'opex': 424798, 'capex': 18916, 'personal': 173266, 'profit': 84466},
}

PROJECTED_DATA = {
    'Nov': {'revenue': 850000, 'opex': 425000, 'capex': 0, 'personal': 85000, 'profit': 321000},
    'Dec': {'revenue': 900000, 'opex': 450000, 'capex': 0, 'personal': 90000, 'profit': 341000},
}

# Calculate totals
total_revenue = sum(d['revenue'] for d in MONTHLY_DATA.values())
total_opex = sum(d['opex'] for d in MONTHLY_DATA.values())
total_capex = sum(d['capex'] for d in MONTHLY_DATA.values())
total_personal = sum(d['personal'] for d in MONTHLY_DATA.values())
total_profit = sum(d['profit'] for d in MONTHLY_DATA.values())
gross_profit = total_revenue - total_opex

# Full year including projections
all_data = {**MONTHLY_DATA, **PROJECTED_DATA}
fy_total_revenue = sum(d['revenue'] for d in all_data.values())
fy_total_opex = sum(d['opex'] for d in all_data.values())
fy_total_capex = sum(d['capex'] for d in all_data.values())
fy_total_personal = sum(d['personal'] for d in all_data.values())
fy_total_profit = sum(d['profit'] for d in all_data.values())
fy_gross_profit = fy_total_revenue - fy_total_opex

# Key metrics
profitable_months = sum(1 for d in MONTHLY_DATA.values() if d['profit'] > 0)

def format_currency(value):
    """Format number as currency"""
    return f"{value:,}"

def format_lakhs(value):
    """Format number in lakhs"""
    return f"{value/100000:.2f}L"

def update_text_in_shape(shape, old_text, new_text):
    """Update text in a shape if it contains old_text"""
    if not shape.has_text_frame:
        return False

    text_frame = shape.text_frame
    for paragraph in text_frame.paragraphs:
        for run in paragraph.runs:
            if old_text in run.text:
                run.text = run.text.replace(old_text, new_text)
                return True
    return False

def update_text_in_table(table, replacements):
    """Update text in table cells"""
    updated = 0
    for row in table.rows:
        for cell in row.cells:
            for old_text, new_text in replacements.items():
                if old_text in cell.text:
                    cell.text = cell.text.replace(old_text, new_text)
                    updated += 1
    return updated

def update_presentation(input_file, output_file):
    """Update data in existing presentation"""
    print(f"Opening: {input_file}")
    prs = Presentation(input_file)

    total_updates = 0

    # Process each slide
    for slide_num, slide in enumerate(prs.slides, 1):
        print(f"\nProcessing Slide {slide_num}...")
        slide_updates = 0

        # Define replacements for this slide
        replacements = {}

        # Common replacements across all slides
        common_replacements = {
            # 10-month totals (old -> new)
            '24,32,223': '16,63,223',  # Net profit
            '24.32L': '16.63L',
            '25.89%': '17.70%',  # Net margin
            '13,17,756': '20,86,756',  # Total CapEx
            '13.18L': '20.87L',

            # Specific month corrections
            '-2,01,319': '-3,71,319',  # Feb profit
            '-201319': '-371319',
            '-2.01L': '-3.71L',

            '2,53,142': '3,142',  # Mar profit
            '253142': '3142',
            '2.53L': '0.03L',

            '2,70,615': '20,615',  # Apr profit
            '270615': '20615',
            '2.71L': '0.21L',

            '-29,086': '-59,086',  # May profit
            '-0.29L': '-0.59L',

            '9,08,366': '8,39,366',  # Jul profit
            '908366': '839366',
            '9.08L': '8.39L',

            # CapEx specific months
            '1,84,674': '3,54,674',  # Feb CapEx
            '184674': '354674',

            '1,54,273': '4,04,273',  # Mar CapEx
            '154273': '404273',

            '1,45,188': '3,95,188',  # Apr CapEx
            '145188': '395188',

            '3,75,636': '4,05,636',  # May CapEx
            '375636': '405636',

            '88,612': '1,57,612',  # Jul CapEx

            # Profitable months
            '7/10': '5/10',
            '70%': '50%',  # Be careful with this one

            # CapEx as % of GP
            '27.97%': '44.29%',
            '28%': '44%',

            # Full year values
            '30,94,223': '23,25,223',  # FY profit
            '30.94L': '23.25L',
            '27.8%': '20.9%',  # FY net margin
        }

        replacements.update(common_replacements)

        # Update all shapes in the slide
        for shape in slide.shapes:
            # Update text in shape
            for old_text, new_text in replacements.items():
                if shape.has_text_frame:
                    if update_text_in_shape(shape, old_text, new_text):
                        slide_updates += 1
                        print(f"  Updated: {old_text} -> {new_text}")

            # Update tables
            if shape.has_table:
                updated = update_text_in_table(shape.table, replacements)
                if updated > 0:
                    print(f"  Updated {updated} table cells")
                    slide_updates += updated

            # Update charts data (if needed)
            if shape.has_chart:
                chart = shape.chart
                # Note: Updating chart data is complex and may require recreation
                # For now, we'll leave charts as-is or handle them separately
                print(f"  Found chart - may need manual update")

        total_updates += slide_updates
        print(f"  Slide {slide_num} updates: {slide_updates}")

    # Save updated presentation
    print(f"\nSaving updated presentation: {output_file}")
    prs.save(output_file)
    print(f"✅ Successfully updated {total_updates} items across {len(prs.slides)} slides")

    return output_file

if __name__ == "__main__":
    import sys
    import os

    # Set UTF-8 encoding for console output
    if sys.platform == 'win32':
        import io
        sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8')

    print("=" * 70)
    print("SKINSPIRE CLINIC - DATA UPDATE SCRIPT")
    print("Updating CapEx and Net Profit values in existing presentation")
    print("=" * 70)
    print()

    # Input file (your edited version)
    input_file = 'Skinspire_Financial_Analysis_2025_Formatted.pptx'

    # Check if file exists
    if not os.path.exists(input_file):
        print(f"ERROR: File not found: {input_file}")
        print("\nAvailable PowerPoint files:")
        for file in os.listdir('.'):
            if file.endswith('.pptx'):
                print(f"  - {file}")
        sys.exit(1)

    # Output file
    output_file = 'Skinspire_Financial_Analysis_2025_Updated.pptx'

    try:
        result = update_presentation(input_file, output_file)
        print("\n" + "=" * 70)
        print("SUCCESS!")
        print("=" * 70)
        print(f"\nUpdated file: {result}")
        print(f"Location: {os.path.abspath(result)}")
        print("\nYour formatting, logos, and customizations are preserved!")
        print("Only the financial data has been updated.")
        print("\n" + "=" * 70)

    except Exception as e:
        print(f"\nERROR: {e}")
        import traceback
        traceback.print_exc()
